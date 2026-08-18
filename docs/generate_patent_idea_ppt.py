#!/usr/bin/env python3
"""专利 IDEA PPT：一种面向智能体工具链的参数引用解析机制"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

W, H = Inches(13.333), Inches(7.5)
OUT = Path(__file__).resolve().parent / "patent_idea_agent_param_ref.pptx"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
TEAL = RGBColor(13, 148, 136)
TEAL_D = RGBColor(15, 118, 110)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(203, 213, 225)
SOFT = RGBColor(240, 253, 250)
BLUE_BG = RGBColor(239, 246, 255)
OK_BG = RGBColor(236, 253, 245)
AMBER_BG = RGBColor(255, 251, 235)
ROSE_BG = RGBColor(255, 241, 242)
CARD = RGBColor(255, 255, 255)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, w=1.25):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def run(p, text, size, bold=False, color=SLATE, font="Microsoft YaHei"):
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(sh, color)
    spTree = slide.shapes._spTree
    el = sh._element
    spTree.remove(el)
    spTree.insert(2, el)


def header(slide, title, sub=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.82))
    fill(bar, NAVY)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.82), W, Inches(0.05))
    fill(acc, TEAL)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.14), Inches(12.5), Inches(0.62))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, title, Pt(18), True, WHITE)
    if sub:
        p2 = tf.add_paragraph()
        run(p2, sub, Pt(11), False, RGBColor(148, 163, 184))


def card(slide, l, t, w, h, color=CARD):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, color)
    stroke(sh, LINE, 1)
    return sh


def tbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def bullets(tf, items, size=Pt(12), color=SLATE, first=True):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        if isinstance(item, tuple):
            run(p, item[0], size, True, TEAL_D)
            run(p, item[1], size, False, color)
        else:
            run(p, "• " + item, size, False, color)
        p.space_after = Pt(4)


def footer(slide, page, total=12):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.28))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run(p, f"专利 IDEA ｜ 一种面向智能体工具链的参数引用解析机制 ｜ {page}/{total}", Pt(9), False, MUTED)


# ---------- 1 Cover ----------
s = blank()
bg(s, NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), W, Inches(2.6))
fill(band, RGBColor(30, 41, 59))
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(0.12), Inches(2.6))
fill(acc, TEAL)
tf = tbox(s, Inches(0.7), Inches(2.4), Inches(12), Inches(2.2))
p = tf.paragraphs[0]
run(p, "专利 IDEA", Pt(14), True, TEAL)
p2 = tf.add_paragraph()
run(p2, "一种面向智能体工具链的参数引用解析机制", Pt(28), True, WHITE)
p3 = tf.add_paragraph()
run(p3, "背景技术 → 技术方案 → 保护点 → 有益效果 → 取证方法 → 检索分析", Pt(13), False, RGBColor(148, 163, 184))
tf2 = tbox(s, Inches(0.7), Inches(5.2), Inches(12), Inches(0.8))
p = tf2.paragraphs[0]
run(p, "模型只表达引用意图　｜　执行前统一拦截求值　｜　真值来自查表/工作记忆　｜　工具零改动", Pt(12), False, RGBColor(203, 213, 225))

# ---------- 2 Agenda ----------
s = blank()
bg(s)
header(s, "目录", "IDEA 撰写结构一览")
items = [
    ("01  ", "来源与公开边界"),
    ("02  ", "背景技术：问题与现有缺陷"),
    ("03  ", "技术方案：架构 / 流程 / 双引用协议"),
    ("04  ", "保护点：方案与架构特征"),
    ("05  ", "有益效果"),
    ("06  ", "取证方法（针对友商）"),
    ("07  ", "专利检索与区别分析"),
    ("08  ", "摘要与待办"),
]
for i, (n, t) in enumerate(items):
    y = Inches(1.15) + Inches(0.68) * i
    c = card(s, Inches(0.6), y, Inches(12.1), Inches(0.58), SOFT if i % 2 == 0 else CARD)
    tf = tbox(s, Inches(0.85), y + Inches(0.12), Inches(11.5), Inches(0.4))
    p = tf.paragraphs[0]
    run(p, n, Pt(16), True, TEAL)
    run(p, t, Pt(15), False, SLATE)
footer(s, 2)

# ---------- 3 Source ----------
s = blank()
bg(s)
header(s, "01  IDEA 来源与公开边界", "产品信息页：识别来源，避免现有技术/公开阻却")
card(s, Inches(0.45), Inches(1.15), Inches(6.1), Inches(5.5), CARD)
tf = tbox(s, Inches(0.7), Inches(1.35), Inches(5.6), Inches(5.1))
p = tf.paragraphs[0]
run(p, "来源说明", Pt(14), True, TEAL_D)
bullets(tf, [
    "来自智能体工具编排产品/项目设计",
    "针对：跨工具传参成本、长值照抄错误",
    "补充场景：手机应用包名等弱语义标识易幻觉",
    "形成统一「参数引用 + 调用前解析」方案",
], Pt(12), SLATE, first=False)

card(s, Inches(6.8), Inches(1.15), Inches(6.0), Inches(5.5), AMBER_BG)
tf = tbox(s, Inches(7.05), Inches(1.35), Inches(5.5), Inches(5.1))
p = tf.paragraphs[0]
run(p, "申报注意（必核）", Pt(14), True, RGBColor(146, 64, 14))
bullets(tf, [
    "已发布产品算现有技术，不能再报 IDEA",
    "已与外界/三方交流过的方案可能构成公开",
    "请在公开前至少 2～3 个月提交 IDEA 电子流",
    "本文未引用外部非公开信息",
    "【待发明人确认】是否已发布 / 是否已对外交流",
], Pt(12), SLATE, first=False)
footer(s, 3)

# ---------- 4 Background ----------
s = blank()
bg(s)
header(s, "02  背景技术", "应用场景 + 现有做法为什么解决不了")
# left scenario
card(s, Inches(0.4), Inches(1.1), Inches(6.2), Inches(2.7), BLUE_BG)
tf = tbox(s, Inches(0.6), Inches(1.25), Inches(5.8), Inches(2.4))
p = tf.paragraphs[0]
run(p, "应用场景", Pt(13), True, RGBColor(30, 64, 175))
bullets(tf, [
    "跨工具传参：上一工具 fileUri / 大字段 → 下一工具入参",
    "弱语义标识：forbidPermission 需要真实 bundleName",
    "包名不规范，用户说「抖音」但工具要 packageName",
], Pt(11), SLATE, first=False)

card(s, Inches(6.8), Inches(1.1), Inches(6.05), Inches(2.7), ROSE_BG)
tf = tbox(s, Inches(7.0), Inches(1.25), Inches(5.65), Inches(2.4))
p = tf.paragraphs[0]
run(p, "现有路径痛点（包名例）", Pt(13), True, RGBColor(159, 18, 57))
bullets(tf, [
    "① 先 getAllInstalledApps 查列表",
    "② 再 forbidPermission(bundleName, …) → 两轮 Loop",
    "凭记忆易编造错误包名；工具多无法逐个改造",
], Pt(11), SLATE, first=False)

card(s, Inches(0.4), Inches(4.0), Inches(12.45), Inches(2.85), CARD)
tf = tbox(s, Inches(0.65), Inches(4.15), Inches(12.0), Inches(2.55))
p = tf.paragraphs[0]
run(p, "现有技术缺陷（简明）", Pt(13), True, TEAL_D)
bullets(tf, [
    "多轮开销：为拿真实标识额外查询；长结果反复回灌抬高 Input Token",
    "幻觉 / 照抄：弱语义标识易编造；长 URI 易截断改写，遵从性差",
    "工具难改 + 真值权威在模型：缺少「真实值只能来自查表/记忆」的结构性约束",
    "业界相近：提示词约束、工作流预置变量替换、密钥引用注入——均难同时满足少轮次、零改工具、消业务标识幻觉",
], Pt(12), SLATE, first=False)
footer(s, 4)

# ---------- 5 Architecture ----------
s = blank()
bg(s)
header(s, "03  技术方案架构", "改进点位置：模型产出 tool call 之后、工具真正执行之前")

boxes = [
    (0.4, 1.3, 2.4, 1.2, "大模型", "输出含 ${…}\n的工具入参", BLUE_BG),
    (3.1, 1.3, 3.2, 1.2, "调度层 + 拦截器", "扫描占位符\n查表求值 / 短路", SOFT),
    (6.6, 1.15, 2.8, 0.7, "工作记忆", "结果引用", AMBER_BG),
    (6.6, 1.95, 2.8, 0.7, "语义词典", "同义词+映射", AMBER_BG),
    (9.7, 1.3, 3.1, 1.2, "工具执行", "只收真值\n实现零改动", OK_BG),
]
for l, t, w, h, title, body, col in boxes:
    card(s, Inches(l), Inches(t), Inches(w), Inches(h), col)
    tf = tbox(s, Inches(l + 0.12), Inches(t + 0.12), Inches(w - 0.2), Inches(h - 0.15))
    p = tf.paragraphs[0]
    run(p, title, Pt(13), True, TEAL_D)
    for line in body.split("\n"):
        p2 = tf.add_paragraph()
        run(p2, line, Pt(11), False, SLATE)

# arrows as text
tf = tbox(s, Inches(2.75), Inches(1.65), Inches(0.4), Inches(0.4))
run(tf.paragraphs[0], "→", Pt(22), True, TEAL)
tf = tbox(s, Inches(9.25), Inches(1.65), Inches(0.4), Inches(0.4))
run(tf.paragraphs[0], "→", Pt(22), True, TEAL)

card(s, Inches(0.4), Inches(2.85), Inches(12.45), Inches(4.15), CARD)
tf = tbox(s, Inches(0.65), Inches(3.05), Inches(12.0), Inches(3.8))
p = tf.paragraphs[0]
run(p, "端到端步骤", Pt(13), True, TEAL_D)
bullets(tf, [
    "S1 工具成功返回 → 默认写入工作记忆，分配短 resultId（对内可映射 toolCallId）",
    "S2 送模前可选掩码：fileUri/fileId 等替换为占位符（所见即所填）",
    "S3 模型填 ${…}，不编造真实包名/长 URI",
    "S4 统一入口拦截：求值成功则回填执行；失败短路，禁止 ${…} 字面量透传工具",
    "S5 语义未命中/歧义 → 可触发查询类工具刷新词典后重试（热路径仍 1 Loop）",
], Pt(12), SLATE, first=False)
footer(s, 5)

# ---------- 6 Dual protocol ----------
s = blank()
bg(s)
header(s, "03  双引用协议 + 语义双表", "同一求值管道，两类数据源")

card(s, Inches(0.4), Inches(1.15), Inches(6.2), Inches(3.5), BLUE_BG)
tf = tbox(s, Inches(0.6), Inches(1.3), Inches(5.8), Inches(3.2))
p = tf.paragraphs[0]
run(p, "A. 结果引用（跨工具传参）", Pt(14), True, RGBColor(30, 64, 175))
bullets(tf, [
    "语法：${resultId.jsonpath}",
    "resultId：字母数字、无「.」；首个「.」切路径",
    "数据源：工作记忆（全量返回优先 → 适配视图）",
    "例：${a1b2c3d4.fileUri}",
], Pt(12), SLATE, first=False)

card(s, Inches(6.8), Inches(1.15), Inches(6.05), Inches(3.5), SOFT)
tf = tbox(s, Inches(7.0), Inches(1.3), Inches(5.65), Inches(3.2))
p = tf.paragraphs[0]
run(p, "B. 语义实体引用（包名等）", Pt(14), True, TEAL_D)
bullets(tf, [
    "语法：${namespace.entity.attribute}",
    "例：${app.抖音.bundleName}",
    "模型只表达「哪个应用」，真值查表",
    "namespace 可扩展：app / contact / device…",
], Pt(12), SLATE, first=False)

card(s, Inches(0.4), Inches(4.85), Inches(12.45), Inches(2.0), AMBER_BG)
tf = tbox(s, Inches(0.65), Inches(5.0), Inches(12.0), Inches(1.7))
p = tf.paragraphs[0]
run(p, "语义词典两张表", Pt(13), True, RGBColor(146, 64, 14))
bullets(tf, [
    "T1 同义词：douyin / Douyin → 抖音（别名归一）",
    "T2 映射：(app, 抖音, bundleName) → com.ss.android.ugc.aweme；0 条失败 / 1 条成功 / 多条歧义",
], Pt(12), SLATE, first=False)
footer(s, 6)

# ---------- 7 Example ----------
s = blank()
bg(s)
header(s, "03  实施例对比", "包名场景：从 2 Loop 到 1 Loop；工具零改动")

card(s, Inches(0.4), Inches(1.15), Inches(6.2), Inches(5.5), ROSE_BG)
tf = tbox(s, Inches(0.6), Inches(1.35), Inches(5.8), Inches(5.1))
p = tf.paragraphs[0]
run(p, "原路径（有缺陷）", Pt(14), True, RGBColor(159, 18, 57))
bullets(tf, [
    "Loop1：getAllInstalledApps",
    "模型从列表抄 bundleName（或凭记忆编造）",
    "Loop2：forbidPermission(真实或错误包名, …)",
    "工具多，无法逐个改造",
    "真值权威仍在模型侧",
], Pt(13), SLATE, first=False)

card(s, Inches(6.8), Inches(1.15), Inches(6.05), Inches(5.5), OK_BG)
tf = tbox(s, Inches(7.0), Inches(1.35), Inches(5.65), Inches(5.1))
p = tf.paragraphs[0]
run(p, "本方案路径", Pt(14), True, RGBColor(21, 128, 61))
bullets(tf, [
    "一次调用业务工具：",
    "forbidPermission(",
    '  bundleName="${app.抖音.bundleName}",',
    '  permission="麦克风")',
    "拦截器：T1 归一 → T2 映射 → 回填真包名",
    "工具无感知；失败才走查询兜底",
], Pt(13), SLATE, first=False)
footer(s, 7)

# ---------- 8 Protection ----------
s = blank()
bg(s)
header(s, "04  保护点", "描述方案/架构/技术实现，不是概念或效果口号")
pts = [
    ("保护点1 · 核心方案", "模型输出统一引用占位符；调度层执行前拦截解析；从工作记忆和/或语义词典回填真值；失败短路不透传；工具无需改造。"),
    ("保护点2 · 双协议", "同时支持 ${resultId.path} 结果引用与 ${ns.entity.attr} 语义引用；同一拦截管道；按命名空间等规则判别类型。"),
    ("保护点3 · 双表消歧", "同义词归一表 + 实体属性映射表，使弱语义业务标识（如包名）真值不由模型生成。"),
    ("保护点4 · 掩码", "送模前将选定字段替换为占位符；模型原样回传；拦截器展开（所见即所填）。"),
    ("保护点5 · 架构", "交互模块 + 带拦截器的调度模块 + 工作记忆 + 语义词典 + 执行模块；拦截器串接于调度与执行之间实现对全工具覆盖。"),
]
for i, (title, body) in enumerate(pts):
    y = Inches(1.05) + Inches(1.1) * i
    card(s, Inches(0.4), y, Inches(12.45), Inches(1.0), SOFT if i % 2 == 0 else CARD)
    tf = tbox(s, Inches(0.65), y + Inches(0.12), Inches(12.0), Inches(0.8))
    p = tf.paragraphs[0]
    run(p, title, Pt(12), True, TEAL_D)
    p2 = tf.add_paragraph()
    run(p2, body, Pt(11), False, SLATE)
footer(s, 8)

# ---------- 9 Benefits ----------
s = blank()
bg(s)
header(s, "05  有益效果", "与背景问题一一对应")
benefits = [
    ("少轮次", "映射命中时业务调用 1 Loop，无需先查全量已安装列表"),
    ("消幻觉", "真实包名/URI 来自词典或工作记忆，模型不编造真值"),
    ("少抄错", "短占位符替代长串复述，降低截断改写；节约 Token/时延"),
    ("零改工具", "统一入口拦截，同类参数工具自动全覆盖"),
    ("可恢复", "失败短路 + 查询兜底刷新后重试，避免错误字面量进工具"),
]
for i, (t, b) in enumerate(benefits):
    col = i % 3
    row = i // 3
    l = 0.4 + col * 4.25
    top = 1.3 + row * 2.7
    card(s, Inches(l), Inches(top), Inches(4.05), Inches(2.4), OK_BG if i % 2 == 0 else BLUE_BG)
    tf = tbox(s, Inches(l + 0.2), Inches(top + 0.35), Inches(3.65), Inches(1.9))
    p = tf.paragraphs[0]
    run(p, t, Pt(16), True, TEAL_D)
    p2 = tf.add_paragraph()
    run(p2, b, Pt(13), False, SLATE)
footer(s, 9)

# ---------- 10 Evidence ----------
s = blank()
bg(s)
header(s, "06  取证方法", "证明友商产品落入保护点；效果表象不能单独作最终证明")
rows = [
    ("引用态→真值态", "对比模型侧 arguments 含 ${…}，实际打到业务接口已是真包名/URI"),
    ("框架统一处理", "多个不同工具均呈同一「引用入参→执行前变真值」模式"),
    ("别名归一", "「抖音」「douyin」首轮不查列表仍落到同一真实包名"),
    ("结果引用", "下一工具实参等于上一工具返回真值，模型参数未见完整抄写"),
    ("失败不透传", "故意非法引用时，工具侧收不到含 ${ 的包名/地址"),
]
for i, (t, b) in enumerate(rows):
    y = Inches(1.1) + Inches(1.05) * i
    card(s, Inches(0.4), y, Inches(2.8), Inches(0.9), SOFT)
    tf = tbox(s, Inches(0.55), y + Inches(0.25), Inches(2.5), Inches(0.55))
    run(tf.paragraphs[0], t, Pt(12), True, TEAL_D)
    card(s, Inches(3.35), y, Inches(9.5), Inches(0.9), CARD)
    tf = tbox(s, Inches(3.55), y + Inches(0.22), Inches(9.1), Inches(0.55))
    run(tf.paragraphs[0], b, Pt(12), False, SLATE)
footer(s, 10)

# Fix footer total - we have more slides. Let me recount and add search + summary slides, then fix page numbers.

# Actually I'll regenerate footers conceptually - we have cover(1) agenda(2) source(3) bg(4) arch(5) dual(6) example(7) protect(8) benefit(9) evidence(10) and need search(11) summary(12). Update footers on previous was wrong for total. Let me add remaining slides with correct numbers and set total=12.

# ---------- 11 Search ----------
s = blank()
bg(s)
header(s, "07  专利检索与区别分析", "区别须落在保护点，功能和效果差异不算专利区别点")

card(s, Inches(0.4), Inches(1.1), Inches(12.45), Inches(1.5), AMBER_BG)
tf = tbox(s, Inches(0.65), Inches(1.25), Inches(12.0), Inches(1.2))
p = tf.paragraphs[0]
run(p, "检索关键词（通用表述）", Pt(12), True, RGBColor(146, 64, 14))
p2 = tf.add_paragraph()
run(p2, "大模型/智能体/工具调用 + 参数 + 占位符/引用/回填 + 执行前/拦截；英文：LLM agent tool calling parameter placeholder resolve before invoke", Pt(11), False, SLATE)

card(s, Inches(0.4), Inches(2.8), Inches(6.2), Inches(4.0), CARD)
tf = tbox(s, Inches(0.6), Inches(2.95), Inches(5.8), Inches(3.7))
p = tf.paragraphs[0]
run(p, "相近文献1：CN120763321A", Pt(12), True, TEAL_D)
bullets(tf, [
    "相近：任务模板占位符由意图信息填充",
    "区别：本发明是工具实参中的引用 + 执行前拦截回填",
    "区别：双协议（结果引用∪语义三段式）+ 双表消歧",
    "区别：失败禁止占位符透传工具",
], Pt(11), SLATE, first=False)

card(s, Inches(6.8), Inches(2.8), Inches(6.05), Inches(4.0), CARD)
tf = tbox(s, Inches(7.0), Inches(2.95), Inches(5.65), Inches(3.7))
p = tf.paragraphs[0]
run(p, "相近公开：密钥/凭证引用注入", Pt(12), True, TEAL_D)
bullets(tf, [
    "相近：执行前将引用解析为真值",
    "区别：对象是业务参数而非密钥；模型仍填参数位",
    "区别：填的是业务引用意图（应用名/结果字段）",
    "区别：目标为消幻觉、少轮次、工具零改覆盖",
], Pt(11), SLATE, first=False)
footer(s, 11)

# ---------- 12 Summary ----------
s = blank()
bg(s)
header(s, "08  摘要与待办", "可直接用于 IDEA 摘要栏")

card(s, Inches(0.4), Inches(1.15), Inches(12.45), Inches(3.6), SOFT)
tf = tbox(s, Inches(0.65), Inches(1.35), Inches(12.0), Inches(3.2))
p = tf.paragraphs[0]
run(p, "摘要", Pt(13), True, TEAL_D)
p2 = tf.add_paragraph()
run(
    p2,
    "本发明公开一种面向智能体工具链的参数引用解析机制。大模型在工具入参中输出统一引用占位符，而不直接生成易幻觉或易抄错的真实标识；工具调度框架在执行前通过统一拦截器解析占位符，分别从会话工作记忆（结果引用）和语义词典（命名空间.实体.属性，经同义词归一与属性映射）查表回填真值后再调用工具，解析失败则短路且不透传字面量。从而在工具零改动前提下，降低多轮查询与长值照抄成本，结构性抑制包名等弱语义标识的幻觉，并适用于跨工具传参与端侧应用操控等场景。",
    Pt(12),
    False,
    SLATE,
)

card(s, Inches(0.4), Inches(4.95), Inches(12.45), Inches(1.9), AMBER_BG)
tf = tbox(s, Inches(0.65), Inches(5.1), Inches(12.0), Inches(1.6))
p = tf.paragraphs[0]
run(p, "发明人待办", Pt(13), True, RGBColor(146, 64, 14))
bullets(tf, [
    "确认未对外发布/未对三方披露；填写产品名与内部案号",
    "按智慧芽指导完成正式检索并替换对比文献",
    "代理人布局权利要求；补绘正式附图后提交电子流",
], Pt(12), SLATE, first=False)
footer(s, 12)

# Fix page footers: cover has no footer; slides 2-12 have footers saying wrong totals on early ones.
# Re-open and patch is hard; regenerate with correct total from the start is cleaner.
# For simplicity, rewrite the script to use total=12 consistently - I already used 10 on early slides.
# Let me just re-run a corrected version by fixing the footer calls - easiest to rewrite file with total=12 throughout.

prs.save(OUT)
print(f"Wrote {OUT}")
