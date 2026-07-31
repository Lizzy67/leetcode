#!/usr/bin/env python3
"""图文并茂：对话内 HTML 渲染方案 PPT（7 节 + 封面）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

ASSETS = "/workspace/docs/assets"
COVER = f"{ASSETS}/ppt-cover-health-chat.png"
SCENARIO = f"{ASSETS}/ppt-scenario-flow.png"
PIPELINE = f"{ASSETS}/ppt-capability-pipeline.png"
SCHEMES = f"{ASSETS}/ppt-three-schemes.png"
SCHEME4 = f"{ASSETS}/ppt-scheme-preview-url.png"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
TEAL = RGBColor(13, 148, 136)
TEAL_D = RGBColor(15, 118, 110)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
SOFT = RGBColor(240, 253, 250)
AMBER = RGBColor(245, 158, 11)
AMBER_BG = RGBColor(255, 251, 235)
OK = RGBColor(4, 120, 87)
GAP = RGBColor(180, 83, 9)
CODE_BG = RGBColor(15, 23, 42)
CODE_FG = RGBColor(203, 213, 225)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, width_pt=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def run(p, text, size, bold=False, color=SLATE, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return r


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(sh, color)
    spTree = slide.shapes._spTree
    el = sh._element
    spTree.remove(el)
    spTree.insert(2, el)


def header(slide, no, title, sub=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.95))
    fill(bar, NAVY)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), W, Inches(0.06))
    fill(acc, TEAL)
    # number pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.22), Inches(0.55), Inches(0.5))
    fill(pill, TEAL)
    pill.adjustments[0] = 0.3
    tf = pill.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], no, Pt(14), True, WHITE)
    tb = slide.shapes.add_textbox(Inches(1.15), Inches(0.18), Inches(11.5), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, title, Pt(22), True, WHITE)
    if sub:
        p2 = tf.add_paragraph()
        run(p2, sub, Pt(11), False, RGBColor(148, 163, 184))


def tbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def round_card(slide, l, t, w, h, fill_c=WHITE, line_c=LINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, fill_c)
    stroke(sh, line_c, 1)
    sh.adjustments[0] = 0.08
    return sh


def arrow_right(slide, l, t, w=Inches(0.35), h=Inches(0.25), color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    fill(sh, color)
    return sh


def set_cell(cell, text, size=9, bold=False, color=SLATE, fill_color=None, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color or WHITE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.clear()
    p.alignment = align
    run(p, text, Pt(size), bold, color)


# ═══════════════════════════════════════════════
# 0 Cover — full-bleed image + overlay
# ═══════════════════════════════════════════════
s = blank()
s.shapes.add_picture(COVER, 0, 0, W, H)
# dark left overlay for title readability
overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(7.2), H)
fill(overlay, NAVY)
# soften: second teal accent bar
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.7), Inches(7.2), Inches(1.8))
fill(band, TEAL_D)
tf = tbox(s, Inches(0.65), Inches(2.0), Inches(6.2), Inches(3.2))
p = tf.paragraphs[0]
run(p, "对话内 HTML 渲染方案", Pt(32), True, WHITE)
p2 = tf.add_paragraph()
run(p2, "Skill 取数 · 图形化生成 · 端侧安全呈现", Pt(15), False, RGBColor(153, 246, 228))
p2.space_before = Pt(14)
p3 = tf.add_paragraph()
run(p3, "以「个人健康分析报告」为验收场景\n图文讲解：问题 → 能力 → 竞品 → Gap → 四方案 → 示例 → 建议", Pt(12), False, RGBColor(148, 163, 184))
p3.space_before = Pt(18)
tf2 = tbox(s, Inches(0.65), Inches(6.05), Inches(6.0), Inches(1.0))
p = tf2.paragraphs[0]
run(p, "技术分享  |  系统方案  |  四路径落地", Pt(13), True, WHITE)

# ═══════════════════════════════════════════════
# 1 问题定义 — 场景图 + 文案
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "01", "问题定义", "需要在端侧对话中直接渲染并显示模型生成的 HTML 内容")

# scenario image
s.shapes.add_picture(SCENARIO, Inches(0.4), Inches(1.25), Inches(7.6), Inches(4.05))
# caption under image
cap = round_card(s, Inches(0.4), Inches(5.4), Inches(7.6), Inches(1.7), SOFT, TEAL)
tf = tbox(s, Inches(0.55), Inches(5.55), Inches(7.3), Inches(1.45))
p = tf.paragraphs[0]
run(p, "场景分镜：取数 → 生成可视化 → 对话内预览", Pt(13), True, TEAL_D)
for line in [
    "① 运动健康 Skill 拉取指标",
    "② 模型生成带图表的 HTML 报告",
    "③ 在会话中直接展示，而非源码/附件",
]:
    p2 = tf.add_paragraph()
    run(p2, line, Pt(12), False, SLATE)
    p2.space_before = Pt(3)

# right cards
round_card(s, Inches(8.2), Inches(1.25), Inches(4.7), Inches(2.55), WHITE, LINE)
tf = tbox(s, Inches(8.4), Inches(1.4), Inches(4.35), Inches(2.3))
p = tf.paragraphs[0]
run(p, "用户原话（示例）", Pt(13), True, TEAL)
p2 = tf.add_paragraph()
run(p2, "「帮我生成最近的个人健康分析报告」", Pt(14), True, NAVY)
p2.space_before = Pt(10)
p3 = tf.add_paragraph()
run(p3, "期望：图表、趋势、结论卡片在对话里可点可看。", Pt(12), False, SLATE)
p3.space_before = Pt(10)

round_card(s, Inches(8.2), Inches(4.0), Inches(4.7), Inches(3.1), AMBER_BG, AMBER)
tf = tbox(s, Inches(8.4), Inches(4.15), Inches(4.35), Inches(2.85))
p = tf.paragraphs[0]
run(p, "成功标准", Pt(13), True, GAP)
for line in [
    "· 会话内直接看见报告",
    "· 图表可交互（不仅截图）",
    "· 无需下载后另开应用",
    "· 安全：脚本不越权主站",
]:
    p2 = tf.add_paragraph()
    run(p2, line, Pt(13), False, SLATE)
    p2.space_before = Pt(6)

# ═══════════════════════════════════════════════
# 2 能力路径 — 配图 + 状态条
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "02", "实现能力路径", "生成 → 落盘 → 同步 → 渲染 → 安全 → 分享")

s.shapes.add_picture(PIPELINE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.55))

# status chips under image aligned to 6 stages
stages = [
    ("1 生成", "LLM + Skill\n产出 HTML", "已有", OK),
    ("2 落盘", "write\nworkspace", "已有", OK),
    ("3 同步", "华为云空间\nREST 可下", "已有", OK),
    ("4 渲染", "iframe /\nWebView", "缺口", GAP),
    ("5 安全", "sandbox\nCSP 隔离", "缺口", GAP),
    ("6 分享", "链接/回放\n转发", "可选", MUTED),
]
for i, (title, desc, st, col) in enumerate(stages):
    x = Inches(0.45) + Inches(i * 2.15)
    card = round_card(s, x, Inches(3.95), Inches(2.05), Inches(2.55), WHITE if st != "缺口" else AMBER_BG, col)
    tf = tbox(s, x + Inches(0.1), Inches(4.1), Inches(1.85), Inches(2.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, title, Pt(13), True, NAVY)
    for line in desc.split("\n"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run(p2, line, Pt(10), False, SLATE)
        p2.space_before = Pt(2)
    # badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.35), Inches(5.95), Inches(1.35), Inches(0.35))
    fill(badge, col)
    badge.adjustments[0] = 0.4
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(btf.paragraphs[0], st, Pt(11), True, WHITE)

tf = tbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.55))
p = tf.paragraphs[0]
run(p, "结论：前三环已通；产品断点在「渲染 + 安全」。健康报告要「在对话里长出来」，必须补齐第 4、5 环。", Pt(13), False, SLATE)

# ═══════════════════════════════════════════════
# 3 竞品表
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "03", "竞品多维对比", "OpenClaw · Cursor · Claude · Gemini · ChatGPT · Manus")

rows = [
    ["维度", "OpenClaw", "Cursor", "Claude", "Gemini", "ChatGPT", "Manus"],
    ["对话内预览", "强\nwidget", "弱*\n偏 IDE", "强\nArtifacts", "中强\nCanvas", "中\nCanvas/ADA", "强\n托管预览"],
    ["呈现形态", "Widget", "工作区文件", "侧栏产物", "Canvas", "Canvas+\n解释器", "站点/App"],
    ["数据/工具", "Skill/工具", "FS/终端", "工具+MCP", "工具+\nWorkspace", "GPTs/ADA", "Agent全流程"],
    ["安全隔离", "双层iframe", "工作区隔离", "独立域CSP", "平台沙箱", "平台沙箱", "托管隔离"],
    ["分享", "pin/渠道", "Git/PR", "Publish", "分享/导出", "分享链接", "一键发布"],
    ["启示", "widget+\n沙箱", "复用write", "单文件\n预览体验", "Canvas\n可参考", "取数可视化", "分享演进"],
]
table = s.shapes.add_table(len(rows), 7, Inches(0.25), Inches(1.2), Inches(12.8), Inches(5.0)).table
widths = [Inches(1.4), Inches(1.9), Inches(1.85), Inches(1.9), Inches(1.9), Inches(1.95), Inches(1.9)]
for i, w in enumerate(widths):
    table.columns[i].width = w
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        if r == 0:
            set_cell(cell, val, 10, True, WHITE, NAVY)
        elif c == 0:
            set_cell(cell, val, 9, True, NAVY, SOFT)
        elif r == 1 and c in (1, 3, 6):
            set_cell(cell, val, 9, True, OK, RGBColor(236, 253, 245))
        else:
            set_cell(cell, val, 9, False, SLATE, WHITE)

tf = tbox(s, Inches(0.35), Inches(6.35), Inches(12.6), Inches(0.9))
p = tf.paragraphs[0]
run(p, "对标优先级：", Pt(12), True, TEAL)
run(p, "Claude Artifacts / OpenClaw Widget  ", Pt(12), False, SLATE)
run(p, "→ 参考 Gemini/ChatGPT Canvas  ", Pt(12), False, SLATE)
run(p, "→ 分享演进看 Manus。", Pt(12), False, SLATE)
p2 = tf.add_paragraph()
run(p2, "* Cursor 强在仓库工程，不是对话内小工具预览主路径。", Pt(10), False, MUTED)
p2.space_before = Pt(4)

# ═══════════════════════════════════════════════
# 4 Gap — 现状/目标可视化流程
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "04", "现状与目标架构 Gap", "复用 write + 云空间；补齐「可预览」闭环")

# AS-IS flow
round_card(s, Inches(0.35), Inches(1.2), Inches(12.6), Inches(2.15), WHITE, LINE)
tf = tbox(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.35))
p = tf.paragraphs[0]
run(p, "现状 AS-IS", Pt(13), True, MUTED)

asis = ["Skill取数", "生成HTML", "write落盘", "云空间同步", "REST下载", "附件/源码?"]
for i, name in enumerate(asis):
    x = Inches(0.55) + Inches(i * 2.05)
    node = round_card(s, x, Inches(1.85), Inches(1.75), Inches(1.15),
                      SOFT if i < 5 else AMBER_BG, TEAL if i < 5 else GAP)
    tf = tbox(s, x + Inches(0.05), Inches(2.15), Inches(1.65), Inches(0.7))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, name, Pt(12), True, NAVY if i < 5 else GAP)
    if i < 5:
        arrow_right(s, x + Inches(1.78), Inches(2.25), Inches(0.28), Inches(0.22), TEAL)

# TO-BE flow
round_card(s, Inches(0.35), Inches(3.5), Inches(12.6), Inches(2.15), SOFT, TEAL)
tf = tbox(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.35))
p = tf.paragraphs[0]
run(p, "目标 TO-BE", Pt(13), True, TEAL_D)

tobe = ["Skill取数", "生成HTML", "write+同步", "发布预览URL", "WebView打开", "会话内交互"]
for i, name in enumerate(tobe):
    x = Inches(0.55) + Inches(i * 2.05)
    node = round_card(s, x, Inches(4.15), Inches(1.75), Inches(1.15), WHITE, TEAL)
    tf = tbox(s, x + Inches(0.05), Inches(4.45), Inches(1.65), Inches(0.7))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, name, Pt(12), True, TEAL_D)
    if i < 5:
        arrow_right(s, x + Inches(1.78), Inches(4.55), Inches(0.28), Inches(0.22), TEAL_D)

# gap chips
gaps = [
    ("P0", "预览容器 WebView/iframe"),
    ("P0", "可打开的预览 URL 或 mime"),
    ("P0", "sandbox / 预览域隔离"),
    ("P1", "同步就绪与刷新降级"),
    ("P2", "流式 / 混合体验增强"),
]
for i, (lv, txt) in enumerate(gaps):
    x = Inches(0.35) + Inches(i * 2.55)
    chip = round_card(s, x, Inches(5.9), Inches(2.45), Inches(1.2),
                      AMBER_BG if lv == "P0" else WHITE, GAP if lv == "P0" else LINE)
    tf = tbox(s, x + Inches(0.12), Inches(6.05), Inches(2.2), Inches(0.95))
    p = tf.paragraphs[0]
    run(p, lv, Pt(12), True, GAP if lv == "P0" else MUTED)
    p2 = tf.add_paragraph()
    run(p2, txt, Pt(11), False, SLATE)

# ═══════════════════════════════════════════════
# 5 四方案 — 2x2 卡片
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "05", "四种实现方案对比", "A 下载后渲染 · B 流式实时 · C 混合 · D 云预览服务出可渲染链接")

schemes = [
    ("方案 1 · A 下载后渲染", "轻量闭环",
     ["write→云空间→REST 拉全文→srcdoc/loadData",
      "稳、贴合现有同步；端要自己注入 HTML",
      "多文件/分享较弱"], TEAL, WHITE),
    ("方案 2 · B 流式实时渲染", "体验增强",
     ["token/delta→端 buffer→边收边刷新",
      "首屏快；半截标签与脚本副作用风险高",
      "需独立推送通道"], MUTED, WHITE),
    ("方案 3 · C 混合", "长报告体验",
     ["流式粗预览 + 完成后落盘定稿再替换",
      "体验近 B，正确性兜底靠完整版",
      "状态机更复杂"], TEAL_D, SOFT),
    ("方案 4 · D 云预览出链接 ★", "中台主路径推荐",
     ["发布到预览 Web 服务→返回 URL→WebView 直接打开",
      "端最简；多文件/分享/统一 CSP 天然友好",
      "需养预览域（或存储签名 URL 网关）"], TEAL, SOFT),
]
for i, (title, badge, lines, col, bgc) in enumerate(schemes):
    r, c = divmod(i, 2)
    x = Inches(0.4) + Inches(c * 6.45)
    y = Inches(1.2) + Inches(r * 3.0)
    round_card(s, x, y, Inches(6.2), Inches(2.8), bgc, col)
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.18), y + Inches(0.18), Inches(2.6), Inches(0.32))
    fill(b, col)
    b.adjustments[0] = 0.4
    btf = b.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(btf.paragraphs[0], badge, Pt(11), True, WHITE)
    tf = tbox(s, x + Inches(0.2), y + Inches(0.6), Inches(5.8), Inches(2.0))
    p = tf.paragraphs[0]
    run(p, title, Pt(15), True, NAVY)
    labels = ["流程", "要点", "代价"]
    for lab, line in zip(labels, lines):
        p2 = tf.add_paragraph()
        run(p2, lab + "：", Pt(12), True, TEAL_D)
        run(p2, line, Pt(12), False, SLATE)
        p2.space_before = Pt(6)

# ═══════════════════════════════════════════════
# 5b 方案4 聚焦
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "05b", "方案 4 详解：云预览服务 → 可渲染链接", "专用（或轻量）HTTP 入口，供 WebView / iframe 直接 loadUrl")

s.shapes.add_picture(SCHEME4, Inches(0.4), Inches(1.2), Inches(7.4), Inches(3.5))

round_card(s, Inches(8.0), Inches(1.2), Inches(4.9), Inches(3.5), SOFT, TEAL)
tf = tbox(s, Inches(8.2), Inches(1.35), Inches(4.55), Inches(3.2))
p = tf.paragraphs[0]
run(p, "核心链路", Pt(14), True, TEAL)
for line in [
    "1. Agent 生成 HTML（可 write 到云空间）",
    "2. 发布到预览 Web 服务 / 签名托管",
    "3. 返回 https://preview…/p/{id}/",
    "4. 端 WebView.loadUrl(link)",
    "5. 服务保证 text/html + CSP/TTL",
]:
    p2 = tf.add_paragraph()
    run(p2, line, Pt(12), False, SLATE)
    p2.space_before = Pt(6)

for i, (t, body) in enumerate([
    ("端侧", "不再必须 REST 拉全文再 srcdoc；\nWeb / 原生统一 loadUrl"),
    ("云侧", "统一承载所有 Web 预览需求；\n多文件相对路径自然支持"),
    ("与云空间关系", "云空间可作源文件仓；\n预览域作渲染入口（可解耦）"),
]):
    x = Inches(0.4) + Inches(i * 4.25)
    round_card(s, x, Inches(4.9), Inches(4.05), Inches(2.2), WHITE, LINE)
    tf = tbox(s, x + Inches(0.2), Inches(5.05), Inches(3.7), Inches(1.9))
    p = tf.paragraphs[0]
    run(p, t, Pt(13), True, TEAL)
    for line in body.split("\n"):
        p2 = tf.add_paragraph()
        run(p2, line, Pt(12), False, SLATE)
        p2.space_before = Pt(4)

# ═══════════════════════════════════════════════
# 6 实现示例
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "06", "实现示例", "健康报告闭环：契约 · 方案 4 主路径 · 方案 A 兜底")

round_card(s, Inches(0.35), Inches(1.2), Inches(4.5), Inches(5.9), WHITE, LINE)
tf = tbox(s, Inches(0.5), Inches(1.35), Inches(4.2), Inches(0.4))
p = tf.paragraphs[0]
run(p, "健康报告串接（方案 4）", Pt(14), True, TEAL)
steps = [
    ("1", "Skill 拉取运动健康数据"),
    ("2", "生成图表化 HTML 报告"),
    ("3", "发布到云预览服务\n（源文件仍可 write 云空间）"),
    ("4", "工具返回 preview_url"),
    ("5", "WebView 直接打开链接"),
]
for i, (n, txt) in enumerate(steps):
    y = Inches(1.9) + Inches(i * 0.95)
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.42), Inches(0.42))
    fill(cir, TEAL)
    ctf = cir.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(ctf.paragraphs[0], n, Pt(12), True, WHITE)
    tf = tbox(s, Inches(1.2), y - Inches(0.05), Inches(3.4), Inches(0.8))
    for j, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        if j:
            p.clear()
        run(p, line, Pt(12), False, SLATE)

round_card(s, Inches(5.05), Inches(1.2), Inches(7.85), Inches(2.55), CODE_BG, NAVY)
tf = tbox(s, Inches(5.25), Inches(1.35), Inches(7.5), Inches(2.25))
code = [
    '// 方案 4：发布后的工具返回（示意）',
    '{ "preview_url": "https://preview.xxx/p/abc/",',
    '  "expires_at": "2026-08-01T00:00:00Z",',
    '  "mime": "text/html",',
    '  "source_file_id": "f_xxx"  // 可选，云空间源文件 }',
]
for i, line in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    if i:
        p.clear()
    run(p, line, Pt(12), False, CODE_FG if i else RGBColor(94, 234, 212), "Consolas")

round_card(s, Inches(5.05), Inches(3.95), Inches(3.8), Inches(3.15), CODE_BG, NAVY)
tf = tbox(s, Inches(5.2), Inches(4.1), Inches(3.5), Inches(2.85))
for i, line in enumerate([
    "// 方案 4 · 端侧",
    "webView.settings.javaScript",
    "  Enabled = true",
    "webView.loadUrl(previewUrl)",
    "// 或 iframe.src = url",
    "// + sandbox=allow-scripts",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    if i:
        p.clear()
    run(p, line, Pt(11), False, CODE_FG if i else RGBColor(94, 234, 212), "Consolas")

round_card(s, Inches(9.05), Inches(3.95), Inches(3.85), Inches(3.15), CODE_BG, NAVY)
tf = tbox(s, Inches(9.2), Inches(4.1), Inches(3.55), Inches(2.85))
for i, line in enumerate([
    "// 方案 A · 兜底",
    "html = await download(id)",
    "iframe.sandbox =",
    "  'allow-scripts'",
    "iframe.srcdoc = html",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    if i:
        p.clear()
    run(p, line, Pt(12), False, CODE_FG if i else RGBColor(94, 234, 212), "Consolas")

# ═══════════════════════════════════════════════
# 7 结论
# ═══════════════════════════════════════════════
s = blank()
bg(s)
header(s, "07", "结论与建议", "方案 4 作 Web 渲染中台；A 作轻量兜底；B/C 作体验增强")

recs = [
    ("01", "问题本质", "对话内安全渲染可视化 HTML\n（取数+生成+呈现），不是再造存储"),
    ("02", "方案 4 定位", "云预览服务生成可渲染链接\nWebView 直接打开——中台主路径"),
    ("03", "与云空间", "云空间作源文件仓；预览域作\nHTTP 渲染入口（可解耦）"),
    ("04", "方案组合", "主路径 D；兜底 A；长等待用 C\n纯 B 不作唯一路径"),
    ("05", "安全底线", "独立预览域 + CSP/TTL\nWebView 仍建议 sandbox 思维"),
    ("06", "下一步", "预览发布 API + 短链 URL\n健康报告验收：会话内 loadUrl"),
]
for i, (no, title, body) in enumerate(recs):
    r, c = divmod(i, 3)
    x = Inches(0.4) + Inches(c * 4.25)
    y = Inches(1.25) + Inches(r * 2.9)
    round_card(s, x, y, Inches(4.05), Inches(2.65), WHITE, TEAL if i < 4 else LINE)
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.5))
    fill(cir, TEAL)
    ctf = cir.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(ctf.paragraphs[0], no, Pt(11), True, WHITE)
    tf = tbox(s, x + Inches(0.85), y + Inches(0.3), Inches(3.0), Inches(0.45))
    p = tf.paragraphs[0]
    run(p, title, Pt(14), True, NAVY)
    tf = tbox(s, x + Inches(0.25), y + Inches(1.0), Inches(3.55), Inches(1.4))
    for j, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        if j:
            p.clear()
        run(p, line, Pt(13), False, SLATE)
        if j:
            p.space_before = Pt(4)

out = "/workspace/docs/agent_html_preview_proposal.pptx"
prs.save(out)
print(f"Saved {out} slides={len(prs.slides)}")
