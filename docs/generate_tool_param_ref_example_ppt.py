#!/usr/bin/env python3
"""One-slide deck for the end-to-end tool-result reference example."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "tool_param_ref_example.pptx"
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(28, 41, 37)
MUTED = RGBColor(107, 122, 116)
FAINT = RGBColor(154, 168, 162)
LINE = RGBColor(221, 229, 225)
TEAL = RGBColor(14, 111, 102)
TEAL_D = RGBColor(11, 84, 78)
WASH = RGBColor(244, 248, 246)
WHITE = RGBColor(255, 255, 255)
MASK = RGBColor(160, 90, 31)
VIS = RGBColor(29, 95, 63)
FONT = "Microsoft YaHei"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
slide = prs.slides.add_slide(prs.slide_layouts[6])


def txt(x, y, w, h, paras, align=PP_ALIGN.LEFT, line_spacing=1.15, mono=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(1)
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = MONO if mono else FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tf


def hline(x, y, w, color=LINE, weight=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(weight))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh


def vline(x, y, h, color=LINE, weight=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Pt(weight), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh


def code_block(x, y, w, h, paras):
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    body.fill.solid(); body.fill.fore_color.rgb = WASH; body.line.fill.background()
    vline(x, y, h, TEAL, 1.6)
    txt(x + 0.1, y + 0.05, w - 0.18, h - 0.1, paras, line_spacing=1.22, mono=True)


def sec_heading(x, y, no, title, note=""):
    txt(x, y, 0.55, 0.24, [[(no, 11, True, TEAL)]])
    txt(x + 0.5, y - 0.01, 3.2, 0.26, [[(title, 12.5, True, INK)]])
    if note:
        txt(x + 3.6, y + 0.03, 9.0, 0.24, [[(note, 8.2, False, FAINT)]])


# background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)

MX = 0.58
CW = 13.333 - 2 * MX

# ---------- header ----------
txt(MX, 0.22, CW, 0.2, [[("A G E N T · T O O L   C A L L I N G · E X A M P L E", 8.5, True, TEAL)]])
txt(MX, 0.42, CW, 0.36, [[("端到端样例：搜图 → 拼图 → 发布", 19, True, INK)]])
txt(MX, 0.83, 12.2, 0.4, [[
    ("用户一句「把上周的滑雪照片拼成一张长图，发到我的相册笔记」，三个工具串联完成。全程模型只写引用：", 9.5, False, INK),
    ("ID 列表和长 URI 被 mask、从未见过真值；长文案可见、供推理，但传参仍是引用。", 9.5, True, TEAL_D),
]], line_spacing=1.28)
hline(MX, 1.32, CW, INK, 1.6)

# ---------- 01 steps ----------
sec_heading(MX, 1.42, "01", "三步执行",
            "字段标记：可见 = 模型能读到内容　mask = 送模前已换成引用，模型只见占位符")

COLW = (CW - 0.72) / 3
CX = [MX, MX + COLW + 0.36, MX + 2 * (COLW + 0.36)]
CY = 1.74

for i in range(1, 3):
    vline(CX[i] - 0.18, CY + 0.02, 3.42)

# ---- step 1 ----
x = CX[0]
txt(x, CY, COLW, 0.16, [[("STEP 1 · 工具 searchPhotos", 7.6, True, FAINT)]])
txt(x, CY + 0.17, COLW, 0.2, [[("搜索照片，结果入库 r1", 10.5, True, TEAL_D)]])
txt(x, CY + 0.4, COLW, 0.14, [[("工具返回（Runtime 所见，resultId = r1）", 7.4, False, FAINT)]])
code_block(x, CY + 0.56, COLW, 0.92, [
    [("{", 7.5, False, INK)],
    [(' "merged_id_list": [563, 728, 991],', 7.5, False, INK)],
    [(' "coverUri": "https://photo.cloud/u9f3a…', 7.5, False, INK)],
    [('   ?sig=eyJhbGciOi…', 7.5, False, INK), ("（214 字符）", 7.5, False, FAINT), ('",', 7.5, False, INK)],
    [(' "summary": "找到3张滑雪照片，1月18日雪场"', 7.5, False, INK)],
    [("}", 7.5, False, INK)],
])
txt(x, CY + 1.54, COLW, 0.14, [[("送模后（模型所见）", 7.4, False, FAINT)]])
code_block(x, CY + 1.7, COLW, 0.78, [
    [(' merged_id_list: ', 7.5, False, INK), ("${r1.merged_id_list}", 7.5, True, TEAL_D)],
    [("   ", 7.5, False, INK), ("← mask·通用", 7.5, True, MASK)],
    [(' coverUri: ', 7.5, False, INK), ("${r1.coverUri}", 7.5, True, TEAL_D), ("  ← mask·通用", 7.5, True, MASK)],
    [(' summary: "找到3张滑雪照片…"  ', 7.5, False, INK), ("← 可见", 7.5, True, VIS)],
])
txt(x, CY + 2.54, COLW, 0.6, [[
    ("模型据 summary 向用户播报「找到3张」；ID 和封面 URI 它不需要懂，只需拿得到。", 8.2, False, MUTED),
]], line_spacing=1.25)

# ---- step 2 ----
x = CX[1]
txt(x, CY, COLW, 0.16, [[("STEP 2 · 工具 createCollage", 7.6, True, FAINT)]])
txt(x, CY + 0.17, COLW, 0.2, [[("引用 ID 列表拼图，结果入库 r2", 10.5, True, TEAL_D)]])
txt(x, CY + 0.4, COLW, 0.14, [[("模型出参（照抄所见占位符）", 7.4, False, FAINT)]])
code_block(x, CY + 0.56, COLW, 0.48, [
    [("createCollage({ images: ", 7.5, False, INK), ("${r1.merged_id_list}", 7.5, True, TEAL_D), (",", 7.5, False, INK)],
    [('   style: "long-strip" })', 7.5, False, INK)],
])
txt(x, CY + 1.1, COLW, 0.14, [[("Runtime beforeHook：先求值，再按 policy 转换", 7.4, False, FAINT)]])
code_block(x, CY + 1.26, COLW, 0.48, [
    [("resolve → ", 7.5, False, FAINT), ("[563, 728, 991]", 7.5, False, INK)],
    [("transform → ", 7.5, False, FAINT), ('[{"file_id":"563"},{"file_id":"728"},…]', 7.5, False, INK)],
])
txt(x, CY + 1.8, COLW, 0.14, [[("工具返回（入库 r2）→ 送模后", 7.4, False, FAINT)]])
code_block(x, CY + 1.96, COLW, 0.78, [
    [(' collageUri: ', 7.5, False, INK), ("${r2.collageUri}", 7.5, True, TEAL_D), ("  ← mask·业务", 7.5, True, MASK)],
    [(' aiDesc: "三张照片自上而下：雪道全景、', 7.5, False, INK)],
    [('   腾空瞬间、雪场合影…', 7.5, False, INK), ("（180 字）", 7.5, False, FAINT), ('"  ', 7.5, False, INK), ("← 可见", 7.5, True, VIS)],
])
txt(x, CY + 2.8, COLW, 0.6, [[
    ("长 URI 被 mask；长文案可见，模型读它确认拼图效果、决定是否直接用作笔记正文。", 8.2, False, MUTED),
]], line_spacing=1.25)

# ---- step 3 ----
x = CX[2]
txt(x, CY, COLW, 0.16, [[("STEP 3 · 工具 publishNote", 7.6, True, FAINT)]])
txt(x, CY + 0.17, COLW, 0.2, [[("可见的用引用传，没见过的照抄", 10.5, True, TEAL_D)]])
txt(x, CY + 0.4, COLW, 0.14, [[("模型出参", 7.4, False, FAINT)]])
code_block(x, CY + 0.56, COLW, 1.06, [
    [("publishNote({", 7.5, False, INK)],
    [(' fileUri: ', 7.5, False, INK), ("${r2.collageUri}", 7.5, True, TEAL_D), ("  场景二：没见过", 7.5, True, MASK)],
    [("   真值，照抄占位符", 7.5, True, MASK)],
    [(' content: ', 7.5, False, INK), ("${r2.aiDesc}", 7.5, True, TEAL_D), ("  场景一：读过内容，", 7.5, True, VIS)],
    [("   传参仍写引用", 7.5, True, VIS)],
    [(' title: "上周滑雪记录" ', 7.5, False, INK), (" 短字面量直接填值 })", 7.5, False, FAINT)],
])
txt(x, CY + 1.68, COLW, 0.14, [[("Runtime 求值后实际调用", 7.4, False, FAINT)]])
code_block(x, CY + 1.84, COLW, 0.78, [
    [(' fileUri: "https://media.cloud/collage/', 7.5, False, INK)],
    [('   7Qx2…?token=…', 7.5, False, INK), ("（198 字符真值）", 7.5, False, FAINT), ('",', 7.5, False, INK)],
    [(' content: "三张照片自上而下：…', 7.5, False, INK), ("（180 字全文）", 7.5, False, FAINT), ('"', 7.5, False, INK)],
])
txt(x, CY + 2.68, COLW, 0.6, [[
    ("发布成功，noteId 再入库。整条链路模型没抄过一个长值，Token 只花在真正需要读的内容上。", 8.2, False, MUTED),
]], line_spacing=1.25)

# ---------- 02 policy ----------
sec_heading(MX, 5.3, "02", "本例的 Policy 配置", "mask 决定「模型看什么」，beforeHook 决定「工具收什么」")
PY_ = 5.6
PCW = (CW - 0.5) / 2
vline(MX + PCW + 0.25, PY_ + 0.02, 1.3)

txt(MX, PY_, PCW, 0.18, [[("Mask Policy：通用兜底 + 业务覆盖", 9.5, True, TEAL_D)]])
code_block(MX, PY_ + 0.22, PCW, 1.06, [
    [("global.maskPolicy ", 7.5, False, INK), ('{ fields: ["*Uri","*Url","*_id_list","fileId"] }', 7.5, False, INK)],
    [("searchPhotos.maskPolicy ", 7.5, False, INK), ('{ keep_visible: ["summary"] }', 7.5, False, INK), ("   # 摘要留给模型", 7.5, False, FAINT)],
    [("createCollage.maskPolicy ", 7.5, False, INK), ('{ keep_visible: ["aiDesc"] }', 7.5, False, INK), ("   # 长文案供推理", 7.5, False, FAINT)],
    [("# coverUri / collageUri / merged_id_list 命中通用规则被 mask；业务规则显式放行可见字段", 7.5, False, FAINT)],
])

x2 = MX + PCW + 0.5
txt(x2, PY_, PCW, 0.18, [[("BeforeHook Policy：入参结构转换", 9.5, True, TEAL_D)]])
code_block(x2, PY_ + 0.22, PCW, 1.06, [
    [("createCollage.beforeHookPolicy {", 7.5, False, INK)],
    [('  images: list(dict("file_id", string(item)) for item in ', 7.5, False, INK), ("$value", 7.5, True, TEAL_D), (")", 7.5, False, INK)],
    [("}", 7.5, False, INK)],
    [("# $value = 引用展开后的真实列表；publishNote 无需转换，仅求值与 Schema 校验", 7.5, False, FAINT)],
])

# ---------- footer ----------
hline(MX, 6.98, CW)
txt(MX, 7.1, CW, 0.28, [[
    ("ID 列表 / 长 URI", 9, True, TEAL_D), ("　mask 后模型从未见过，占位符原样穿过三步　　", 8.2, False, MUTED),
    ("长文案 aiDesc", 9, True, TEAL_D), ("　可见、可推理，传参仍是 ${r2.aiDesc}　　", 8.2, False, MUTED),
    ("短字面量 title", 9, True, TEAL_D), ("　模型自己生成，照常直接填值", 8.2, False, MUTED),
]], PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Wrote {OUT}")
