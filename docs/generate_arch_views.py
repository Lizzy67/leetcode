#!/usr/bin/env python3
"""独立 PPT：逻辑架构图 + 运行视图（按已确认口径）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
TEAL = RGBColor(13, 148, 136)
TEAL_D = RGBColor(15, 118, 110)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(203, 213, 225)
SOFT = RGBColor(240, 253, 250)
BLUE_BG = RGBColor(239, 246, 255)
OK_BG = RGBColor(236, 253, 245)
AMBER_BG = RGBColor(255, 251, 235)
GAP = RGBColor(180, 83, 9)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, w=1.25):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def run(p, text, size, bold=False, color=SLATE, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(sh, color)
    spTree = slide.shapes._spTree
    el = sh._element
    spTree.remove(el)
    spTree.insert(2, el)


def header(slide, title, sub):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.85))
    fill(bar, NAVY)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), W, Inches(0.05))
    fill(acc, TEAL)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run(p, title, Pt(20), True, WHITE)
    p2 = tf.add_paragraph()
    run(p2, sub, Pt(11), False, RGBColor(148, 163, 184))


def tbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def card(slide, l, t, w, h, title, lines, fill_c=WHITE, line_c=TEAL, title_c=TEAL_D):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, fill_c)
    stroke(sh, line_c)
    sh.adjustments[0] = 0.08
    tf = tbox(slide, l + Inches(0.12), t + Inches(0.1), w - Inches(0.22), h - Inches(0.16))
    p = tf.paragraphs[0]
    run(p, title, Pt(12), True, title_c)
    for line in lines:
        p2 = tf.add_paragraph()
        run(p2, line, Pt(11), False, SLATE)
        p2.space_before = Pt(3)
    return sh


def lane(slide, l, t, w, h, label, fill_c):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, fill_c)
    stroke(sh, LINE)
    sh.adjustments[0] = 0.03
    tf = tbox(slide, l + Inches(0.12), t + Inches(0.08), w - Inches(0.2), Inches(0.3))
    p = tf.paragraphs[0]
    run(p, label, Pt(11), True, MUTED)
    return sh


# ───────── Cover notes slide optional? Keep 3 slides: cover brief + logic + runtime ─────────
# User asked for architecture pages; 2 content slides + 1 decision note is helpful.

s = blank()
bg(s)
header(s, "架构视图（独立稿）", "不并入主 PPT · 已按评审口径收敛：独立预览服务 / 直发预览 / 单文件 / WebView")

decisions = [
    ("预览服务", "独立新服务（Preview Web Service）"),
    ("发布路径", "HTML 直发预览服务；云空间仅备份"),
    ("一期范围", "只做单文件 HTML（CSS/JS 内联）"),
    ("端侧呈现", "WebView 组件直接打开 preview_url"),
    ("鉴权", "不引入 AK/SK；publish 后签发短时可访问 preview_url"),
    ("不进本图", "流式方案 B/C（见备注，非本期主架构）"),
]
for i, (k, v) in enumerate(decisions):
    r, c = divmod(i, 3)
    x = Inches(0.4) + Inches(c * 4.25)
    y = Inches(1.3) + Inches(r * 2.7)
    card(s, x, y, Inches(4.05), Inches(2.4), k, [v], SOFT if i < 4 else AMBER_BG, TEAL if i < 4 else GAP, TEAL_D if i < 4 else GAP)

# ───────── Logical ─────────
s = blank()
bg(s)
header(s, "逻辑架构图", "方案4主路径：Agent 直发独立预览服务 → 返回 URL → 端 WebView 渲染；云空间仅备份")

# Client
lane(s, Inches(0.3), Inches(1.05), Inches(12.7), Inches(1.45), "① 端侧", BLUE_BG)
card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(0.9), "对话 UI", ["发起任务；展示预览入口"], WHITE, RGBColor(37, 99, 235), RGBColor(37, 99, 235))
card(s, Inches(6.7), Inches(1.4), Inches(5.95), Inches(0.9), "WebView 组件 ★", ["loadUrl(preview_url) 会话内渲染"], WHITE, TEAL)

# Agent
lane(s, Inches(0.3), Inches(2.65), Inches(12.7), Inches(1.7), "② Agent 运行时", SOFT)
card(s, Inches(0.55), Inches(3.05), Inches(2.95), Inches(1.1), "编排 / Session", ["对话上下文", "工具调用编排"], WHITE, TEAL)
card(s, Inches(3.65), Inches(3.05), Inches(2.95), Inches(1.1), "LLM", ["生成单文件 HTML", "图表/文案内联"], WHITE, TEAL)
card(s, Inches(6.75), Inches(3.05), Inches(2.95), Inches(1.1), "Skills", ["运动健康数据 Skill", "其它取数…"], WHITE, TEAL)
card(s, Inches(9.85), Inches(3.05), Inches(2.85), Inches(1.1), "Tools", ["publish_preview ★", "write（备份云空间）"], WHITE, TEAL_D)

# Preview primary + backup secondary
lane(s, Inches(0.3), Inches(4.5), Inches(7.4), Inches(2.55), "③ Web 渲染中台（独立新服务）★", SOFT)
card(s, Inches(0.5), Inches(4.95), Inches(6.95), Inches(1.85), "Preview Web Service（独立部署）", [
    "接收：Agent 直发的单文件 HTML",
    "输出：可访问 preview_url",
    "职责：托管、text/html、CSP、TTL、访问控制",
], WHITE, TEAL_D)

lane(s, Inches(7.9), Inches(4.5), Inches(5.1), Inches(2.55), "④ 备份（非渲染入口）", OK_BG)
card(s, Inches(8.1), Inches(4.95), Inches(4.7), Inches(1.85), "华为云空间", [
    "write 同步仅作备份/审计",
    "不承担 WebView 渲染入口",
    "需要时可 REST 取源文件",
], WHITE, TEAL)

tf = tbox(s, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.25))
p = tf.paragraphs[0]
run(p, "鉴权注记：不引入 AK/SK；预览服务签发短时 preview_url，WebView 仅打开该链接", Pt(11), False, MUTED)

# ───────── Runtime ─────────
s = blank()
bg(s)
header(s, "运行视图 · 健康报告", "直发预览服务为主；云空间备份并行/事后；端仅 WebView 打开链接")

cols = ["用户 / 端侧", "Agent 运行时", "Skill", "预览 Web 服务", "云空间(备份)"]
col_x = [Inches(0.35), Inches(2.9), Inches(5.45), Inches(8.0), Inches(10.55)]
for i, name in enumerate(cols):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x[i], Inches(1.05), Inches(2.35), Inches(0.45))
    fill(sh, NAVY)
    sh.adjustments[0] = 0.2
    tf = sh.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], name, Pt(11), True, WHITE)
    spine = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[i] + Inches(1.12), Inches(1.55), Inches(0.035), Inches(5.45))
    fill(spine, LINE)

events = [
    (1.7, 0, "1. 请求健康报告", TEAL),
    (2.25, 1, "2. 编排任务", TEAL),
    (2.8, 2, "3. 取运动健康数据", TEAL),
    (3.35, 1, "4. LLM 生成单文件 HTML", TEAL),
    (3.9, 3, "5. publish_preview 直发", TEAL_D),
    (4.45, 3, "6. 返回 preview_url", TEAL_D),
    (5.0, 4, "7. write 备份源文件", MUTED),
    (5.55, 1, "8. 回传 URL 给端", TEAL),
    (6.1, 0, "9. WebView 打开链接", TEAL),
    (6.65, 3, "10. 返回 HTML 完成渲染", TEAL),
]
for y, col, text, color in events:
    x = col_x[col]
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.08), Inches(y), Inches(2.2), Inches(0.42))
    fill(sh, color)
    sh.adjustments[0] = 0.2
    tf = sh.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], text, Pt(10), True, WHITE)

leg = tbox(s, Inches(0.35), Inches(7.15), Inches(12.6), Inches(0.25))
p = leg.paragraphs[0]
run(p, "青色=主路径　　灰色=备份路径（可与发布并行，不阻塞预览）　　WebView 仅打开短时 preview_url，不携带 AK/SK", Pt(10), False, MUTED)

# ───────── Auth note slide ─────────
s = blank()
bg(s)
header(s, "说明：鉴权口径 · 单文件为何够用 · 流式 B/C 是什么", "供对齐，不并入主 PPT")

card(s, Inches(0.35), Inches(1.2), Inches(4.15), Inches(5.8), "鉴权口径（已定）", [
    "不引入 AK/SK。",
    "",
    "Agent 调用预览服务：",
    "  走现有服务间可信通道",
    "  （内网 / 会话态等，不落 AKSK）",
    "",
    "端侧 WebView：",
    "  只打开返回的 preview_url",
    "  URL 自身短时有效即可",
    "",
    "原则：",
    "  密钥不进 WebView",
    "  不进 HTML，不进客户端本地",
], SOFT, TEAL)

card(s, Inches(4.65), Inches(1.2), Inches(4.15), Inches(5.8), "静态目录资源有啥用？", [
    "指一个页面拆成：",
    "  index.html + app.css + chart.js",
    "  + 图片等多个文件",
    "",
    "预览服务若支持「目录托管」，",
    "相对路径才能一起打开。",
    "",
    "一期单文件：",
    "  CSS/JS 都写进同一个 HTML",
    "  → 不需要目录能力",
    "",
    "所以一期不做静态目录完全合理；",
    "以后报告变复杂再加即可。",
], WHITE, TEAL)

card(s, Inches(8.95), Inches(1.2), Inches(4.0), Inches(5.8), "流式 B/C 是什么？", [
    "这是主 PPT 里另外两条",
    "「端内喂 HTML」的备选，",
    "不是方案4的一部分：",
    "",
    "B：模型边生成，边把 HTML",
    "   片段推到端上刷新",
    "",
    "C：先流式粗看，完成后再",
    "   落盘/定稿替换",
    "",
    "你们定主路径=预览服务出链接后，",
    "B/C 可先不进架构图。",
    "",
    "本独立 PPT 已按此省略。",
], AMBER_BG, GAP, GAP)

out = "/workspace/docs/agent_html_arch_views.pptx"
prs.save(out)
print(f"Saved {out} slides={len(prs.slides)}")
