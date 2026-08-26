#!/usr/bin/env python3
"""One-slide deck for the tool-result reference scheme, editorial style."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "tool_param_ref_onepage.pptx"
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(28, 41, 37)
MUTED = RGBColor(107, 122, 116)
FAINT = RGBColor(154, 168, 162)
LINE = RGBColor(221, 229, 225)
TEAL = RGBColor(14, 111, 102)
TEAL_D = RGBColor(11, 84, 78)
WASH = RGBColor(244, 248, 246)
WHITE = RGBColor(255, 255, 255)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
slide = prs.slides.add_slide(prs.slide_layouts[6])


def txt(x, y, w, h, runs, align=PP_ALIGN.LEFT, line_spacing=1.15):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tf


def hline(x, y, w, color=LINE, weight=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(weight))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def vline(x, y, h, color=LINE, weight=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Pt(weight), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def sec_heading(x, y, no, title, note=""):
    txt(x, y, 0.55, 0.24, [[(no, 11, True, TEAL)]])
    txt(x + 0.52, y - 0.015, 3.4, 0.26, [[(title, 12.5, True, INK)]])
    if note:
        txt(x + 3.9, y + 0.02, 8.6, 0.24, [[(note, 8.5, False, FAINT)]])


# background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)

MX = 0.62
CW = 13.333 - 2 * MX

# ---------- header ----------
txt(MX, 0.26, CW, 0.2, [[("A G E N T · T O O L   C A L L I N G", 9, True, TEAL)]])
txt(MX, 0.47, CW, 0.4, [[("工具结果引用：让模型写「引用」，让 Runtime 填「真值」", 20.5, True, INK)]])
txt(MX, 0.9, 12.0, 0.42, [[
    ("上一个工具的结果不再由模型抄进下游参数，而是写一个引用 ${resultId.path}；Runtime 在工具执行前从工作记忆取出真值、按需转换、校验合法，再去调用工具。", 10, False, INK),
    [("真值不经模型之手，存量工具一行不改。", 10, True, TEAL_D)][0] if False else ("", 1, False, INK),
]])
txt(MX, 1.3, 11.6, 0.22, [[("真值不经模型之手，存量工具一行不改。", 10, True, TEAL_D)]])
hline(MX, 1.58, CW, INK, 1.6)

# ---------- 01 why ----------
sec_heading(MX, 1.7, "01", "为什么要做")
txt(MX, 1.96, CW, 0.42, [[
    ("下一个工具的参数常来自上一个工具的结果：文件 URI、对象 ID、大段文本、一整个列表。现在要先整段回灌上下文，再由模型原样抄进下游参数——长 URI 抄着抄着就断了；结构不一致时还要模型改写大段 JSON，每一轮都在烧 Token。", 9.5, False, INK),
]], line_spacing=1.25)
txt(MX, 2.42, CW, 0.2, [[
    ("抄不准 ", 9, True, INK), ("长值复述易截断改写　　", 9, False, MUTED),
    ("太昂贵 ", 9, True, INK), ("大结果反复回灌，Token与时延一起涨　　", 9, False, MUTED),
    ("改不动 ", 9, True, INK), ("结构适配靠模型改JSON或逐个改工具，都不划算", 9, False, MUTED),
]])

# ---------- 02 pipeline ----------
sec_heading(MX, 2.76, "02", "方案主链路", "Reference → Resolve → Invoke，全部发生在工具执行之前")
nodes = [
    ("工具 A", "结果入库", "结果默认写入工作记忆，得到短标识 resultId"),
    ("Runtime", "按 Policy 送模", "命中 mask policy 的字段换成引用，其余全量可见"),
    ("模型", "只写引用", "下游参数填 ${resultId.path}，可见与否都不复述真值"),
    ("Runtime", "求值·转换·校验", "查记忆回填真值，按策略适配结构，验完Schema与权限"),
    ("工具 B", "真值执行", "校验通过才调用；结果再入库，链路闭环"),
]
NY, NH = 3.04, 0.92
NW, GAP = 2.22, 0.24
for i, (who, what, how) in enumerate(nodes):
    nx = MX + i * (NW + GAP)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(nx), Inches(NY), Inches(NW), Pt(2.2))
    top.fill.solid(); top.fill.fore_color.rgb = TEAL; top.line.fill.background()
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(nx), Inches(NY + 0.03), Inches(NW), Inches(NH))
    body.fill.solid(); body.fill.fore_color.rgb = WASH; body.line.fill.background()
    txt(nx + 0.1, NY + 0.09, NW - 0.2, 0.16, [[(who, 7.6, True, FAINT)]])
    txt(nx + 0.1, NY + 0.26, NW - 0.2, 0.19, [[(what, 10, True, INK)]])
    txt(nx + 0.1, NY + 0.47, NW - 0.2, 0.44, [[(how, 8, False, MUTED)]], line_spacing=1.18)
    if i < len(nodes) - 1:
        txt(nx + NW + 0.015, NY + 0.33, 0.22, 0.26, [[("→", 12, True, TEAL)]], PP_ALIGN.CENTER)

# ---------- 03 two scenarios ----------
sec_heading(MX, 4.12, "03", "两种场景", "按「模型要不要看到这个值」分流，殊途同归于引用")
SC_Y = 4.4
SCW = (CW - 0.5) / 2

# scenario 1
txt(MX, SC_Y, SCW, 0.18, [[("场景一 · 可见但引用", 8.2, True, FAINT)]])
txt(MX, SC_Y + 0.19, SCW, 0.2, [[("结果全量进上下文，参数仍写引用", 10.5, True, TEAL_D)]])
txt(MX, SC_Y + 0.43, SCW, 0.56, [[
    ("模型需要读内容才能决策：工具返回原样进入上下文，推理照常。唯一约束是传给下游时参数写 ${resultId.path}，不许把看过的内容再抄一遍——「看」与「传」分开，由 Skill 约定与提示词共同约束。", 8.4, False, INK),
]], line_spacing=1.22)
m1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MX), Inches(SC_Y + 1.02), Inches(SCW), Inches(0.42))
m1.fill.solid(); m1.fill.fore_color.rgb = WASH; m1.line.fill.background()
vline(MX, SC_Y + 1.02, 0.42, TEAL, 1.6)
txt(MX + 0.12, SC_Y + 1.06, SCW - 0.2, 0.36, [
    [("上下文：", 7.8, False, FAINT), ('text = "血糖偏高，建议…"（全量可见）', 7.8, False, INK)],
    [("模型出参：", 7.8, False, FAINT), ("report(content = ${e5f6a7b8.text})", 7.8, True, TEAL_D)],
], line_spacing=1.25)

# scenario 2
S2X = MX + SCW + 0.5
vline(S2X - 0.25, SC_Y + 0.05, 1.42)
txt(S2X, SC_Y, SCW, 0.18, [[("场景二 · Mask 后不可见", 8.2, True, FAINT)]])
txt(S2X, SC_Y + 0.19, SCW, 0.2, [[("Policy 注入掩码，模型全程不接触真值", 10.5, True, TEAL_D)]])
txt(S2X, SC_Y + 0.43, SCW, 0.56, [[
    ("模型不需要理解的值（文件URI、对象ID等指针字段），送模前就替换成引用。规则由两层 Policy 注入：通用 Policy 按字段名兜底（fileUri、fileId、*Url…），业务 Policy 按工具粒度增删覆盖，业务优先。", 8.4, False, INK),
]], line_spacing=1.22)
m2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(S2X), Inches(SC_Y + 1.02), Inches(SCW), Inches(0.42))
m2.fill.solid(); m2.fill.fore_color.rgb = WASH; m2.line.fill.background()
vline(S2X, SC_Y + 1.02, 0.42, TEAL, 1.6)
txt(S2X + 0.12, SC_Y + 1.06, SCW - 0.2, 0.36, [
    [("真值：", 7.8, False, FAINT), ('fileUri = "https://…/very/long/path?token=…"', 7.8, False, INK)],
    [("送模后：", 7.8, False, FAINT), ("fileUri = ${a1b2c3d4.fileUri}", 7.8, True, TEAL_D), ("  ← 所见即所填", 7.8, False, MUTED)],
], line_spacing=1.25)

# ---------- 04 guarantees ----------
sec_heading(MX, 5.98, "04", "真值与保障")
GY = 6.24
txt(MX, GY, SCW, 0.6, [
    [("真值从哪来、怎么变  ", 9, True, TEAL_D),
     ("只从工作记忆取（全量结果集→工具适配视图）；resultId 由 Runtime 生成并映射 toolCallId。结构不一致时按参数策略转换：[563,728] → [{\"file_id\":\"563\"},…]。", 8.4, False, INK)],
], line_spacing=1.25)
vline(S2X - 0.25, GY + 0.02, 0.52)
txt(S2X, GY, SCW, 0.6, [
    [("出错怎么办  ", 9, True, TEAL_D),
     ("解析失败、字段缺失、转换出错一律短路，${…} 绝不透传，错误回模型改写重试；结果带作用域与有效期，并行链路必须显式写 resultId。", 8.4, False, INK)],
], line_spacing=1.25)

# ---------- footer ----------
hline(MX, 6.9, CW)
txt(MX, 7.02, CW, 0.3, [[
    ("省 Token", 9.5, True, TEAL_D), ("　大结果不再回灌复述　　", 8.5, False, MUTED),
    ("抄不错", 9.5, True, TEAL_D), ("　真值不经模型之手　　", 8.5, False, MUTED),
    ("结构解耦", 9.5, True, TEAL_D), ("　Schema差异由策略消化　　", 8.5, False, MUTED),
    ("零改动", 9.5, True, TEAL_D), ("　存量工具无感接入", 8.5, False, MUTED),
]], PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Wrote {OUT}")
