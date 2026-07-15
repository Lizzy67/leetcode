#!/usr/bin/env python3
"""Generate architecture PPT: app memory × Xiaoyi integration schemes."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree

C_DARK = RGBColor(0x12, 0x1A, 0x22)
C_CARD = RGBColor(0xF4, 0xF7, 0xF9)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_INK = RGBColor(0x1A, 0x2B, 0x3C)
C_MUTED = RGBColor(0x5A, 0x6B, 0x7C)
C_PRIMARY = RGBColor(0x0E, 0x7C, 0x86)
C_PRIMARY_D = RGBColor(0x0A, 0x5C, 0x64)
C_INTERNAL = RGBColor(0x2F, 0x6F, 0xE5)
C_EXTERNAL = RGBColor(0x0E, 0x7C, 0x86)
C_MEMORY = RGBColor(0xC4, 0x5C, 0x26)
C_LINE = RGBColor(0xD0, 0xD8, 0xE0)
C_SOFT = RGBColor(0xE8, 0xF0, 0xF3)
C_SOFT_B = RGBColor(0xE3, 0xEC, 0xFB)
C_SOFT_T = RGBColor(0xD9, 0xEF, 0xF1)
C_SOFT_M = RGBColor(0xFA, 0xEB, 0xE0)
C_WARN = RGBColor(0xB5, 0x6E, 0x1A)

TOTAL = 12
FONT = "Microsoft YaHei"


def set_run(run, size=18, bold=False, color=C_INK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = etree.SubElement(
            r_pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea"
        )
    ea.set("typeface", font)


def add_text(shape, text, size=18, bold=False, color=C_INK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return tf


def add_para(tf, text, size=14, bold=False, color=C_INK, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    if radius:
        s.adjustments[0] = 0.08
    return s


def box(
    slide,
    x,
    y,
    w,
    h,
    title,
    subtitle=None,
    fill=C_CARD,
    title_color=C_INK,
    sub_color=C_MUTED,
    title_size=14,
    sub_size=11,
    border=None,
):
    s = rect(slide, x, y, w, h, fill, line=border, radius=True)
    tf = s.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(6)
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, size=title_size, bold=True, color=title_color)
    if subtitle:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = subtitle
        set_run(r, size=sub_size, bold=False, color=sub_color)
    return s


def arrow_right(slide, x, y, w=0.45, h=0.22, color=C_PRIMARY):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def arrow_down(slide, x, y, w=0.22, h=0.4, color=C_PRIMARY):
    s = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def label(slide, x, y, w, h, text, size=12, bold=False, color=C_MUTED, align=PP_ALIGN.CENTER):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    add_text(s, text, size=size, bold=bold, color=color, align=align)
    return s


def title_bar(slide, title, subtitle=None, page=None):
    rect(slide, 0, 0, 13.333, 0.08, C_PRIMARY)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(10.5), Inches(0.5))
    add_text(t, title, size=28, bold=True, color=C_INK)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.35))
        add_text(s, subtitle, size=14, color=C_MUTED)
    if page is not None:
        p = slide.shapes.add_textbox(Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3))
        add_text(p, f"{page}/{TOTAL}", size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)
    rect(slide, 0.6, 7.0, 12.1, 0.01, C_LINE)


def section_chip(slide, x, y, text, fill):
    s = rect(slide, x, y, 1.7, 0.32, fill, radius=True)
    add_text(s, text, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    return s


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, C_WHITE)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, C_DARK)
    rect(slide, 0, 0, 0.18, 7.5, C_PRIMARY)
    rect(slide, 0, 6.9, 13.333, 0.6, C_PRIMARY_D)
    add_text(
        slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11), Inches(1)),
        "应用记忆数据 × 小艺对接方案",
        size=36,
        bold=True,
        color=C_WHITE,
    )
    add_text(
        slide.shapes.add_textbox(Inches(0.9), Inches(3.1), Inches(11), Inches(0.6)),
        "内部记忆管理 vs 外部记忆共享｜含小艺记忆 / 帮记路径",
        size=18,
        color=RGBColor(0xB8, 0xC8, 0xD0),
    )
    add_text(
        slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11), Inches(0.4)),
        "架构方案对比 · 数据流 · 落地建议",
        size=14,
        color=C_PRIMARY,
    )
    add_text(
        slide.shapes.add_textbox(Inches(0.9), Inches(7.05), Inches(10), Inches(0.3)),
        "HarmonyOS / 小艺开放能力梳理",
        size=12,
        color=RGBColor(0xD0, 0xE8, 0xEA),
    )

    # 2 TOC
    slide = blank(prs)
    title_bar(slide, "目录", "先分清两类记忆，再看四条可落地架构", 2)
    items = [
        ("01", "问题定义与边界", "应用记忆如何被小艺使用"),
        ("02", "概念分层", "内部记忆管理 / 外部记忆共享 / 小艺记忆"),
        ("03", "总览架构", "一张图看清数据权属与流向"),
        ("04", "方案 A", "Intents Kit 外部意图共享"),
        ("05", "方案 A 数据流", "共享与调用双向闭环"),
        ("06", "方案 B", "应用内 / Agent 长期记忆（内部）"),
        ("07", "方案 C", "小艺记忆 / 帮记 + 记忆授权"),
        ("08", "方案 D", "双轨组合架构（推荐）"),
        ("09", "对比与落地", "选型建议与接入步骤"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col, row = divmod(i, 5)
        if i >= 5:
            col, row = 1, i - 5
        else:
            col, row = 0, i
        x = 0.7 + col * 6.2
        y = 1.3 + row * 1.05
        rect(slide, x, y, 5.8, 0.9, C_SOFT if col == 0 else C_SOFT_T, radius=True)
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(0.8), Inches(0.45)),
            num,
            size=20,
            bold=True,
            color=C_PRIMARY,
        )
        add_text(
            slide.shapes.add_textbox(Inches(x + 1.1), Inches(y + 0.15), Inches(4.4), Inches(0.3)),
            title,
            size=15,
            bold=True,
            color=C_INK,
        )
        add_text(
            slide.shapes.add_textbox(Inches(x + 1.1), Inches(y + 0.48), Inches(4.4), Inches(0.3)),
            desc,
            size=12,
            color=C_MUTED,
        )

    # 3 Problem
    slide = blank(prs)
    title_bar(slide, "问题定义与边界", "「应用记忆」≠「小艺记忆」；对接目标决定方案选择", 3)
    box(
        slide,
        0.6,
        1.35,
        3.9,
        2.4,
        "应用记忆数据",
        "偏好 · 行为轨迹 · 业务实体\n会话上下文 · 用户画像片段",
        fill=C_SOFT_B,
        title_color=C_INTERNAL,
        title_size=16,
        sub_size=12,
    )
    box(
        slide,
        4.85,
        1.9,
        3.5,
        1.3,
        "要解决什么？",
        "让系统级语音助手（小艺）\n理解并使用这些数据",
        fill=C_SOFT,
        title_color=C_PRIMARY,
        title_size=15,
        sub_size=12,
    )
    box(
        slide,
        8.7,
        1.35,
        3.9,
        2.4,
        "小艺侧能力",
        "对话理解 · 搜索 · 建议\n小艺记忆/帮记 · 智能体编排",
        fill=C_SOFT_M,
        title_color=C_MEMORY,
        title_size=16,
        sub_size=12,
    )
    arrow_right(slide, 4.55, 2.35, 0.28, 0.22, C_MUTED)
    arrow_right(slide, 8.4, 2.35, 0.28, 0.22, C_MUTED)

    goals = [
        ("G1 可被调用", "用户对小艺说话，能触发应用能力", "主通道：意图调用", C_PRIMARY),
        ("G2 可被推荐", "系统学习习惯后主动建议", "主通道：意图共享", C_INTERNAL),
        ("G3 可被记忆", "长期偏好进入助手记忆体系", "系统记忆为主，三方受限", C_MEMORY),
    ]
    for i, (t, d, note, c) in enumerate(goals):
        x = 0.6 + i * 4.2
        rect(slide, x, 4.15, 3.95, 2.3, C_CARD, line=c, radius=True)
        chip = rect(slide, x + 0.2, 4.35, 1.55, 0.32, c, radius=True)
        add_text(chip, t, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        body = slide.shapes.add_textbox(Inches(x + 0.25), Inches(4.9), Inches(3.45), Inches(1.3))
        tf = add_text(body, d, size=14, color=C_INK)
        add_para(tf, note, size=12, color=C_MUTED if i < 2 else C_WARN, space_before=10)

    # 4 Concepts
    slide = blank(prs)
    title_bar(
        slide,
        "概念分层：内部记忆 vs 外部记忆共享",
        "先定权属，再定接口；避免把「应用本地记忆」和「系统助手记忆」混为一谈",
        4,
    )
    rect(slide, 0.55, 1.35, 6.0, 5.2, C_SOFT_B, radius=True)
    section_chip(slide, 0.8, 1.55, "内部记忆管理", C_INTERNAL)
    add_text(
        slide.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(5.5), Inches(0.4)),
        "数据留在应用 / 智能体域内",
        size=15,
        bold=True,
        color=C_INK,
    )
    for i, p in enumerate(
        [
            "应用自建 Memory Center（情景/事实/程序）",
            "小艺开放平台智能体「长期记忆」编排",
            "仅服务本应用对话与个性化，不默认进系统",
            "隐私边界清晰，可控可删可版本化",
        ]
    ):
        box(slide, 0.85, 2.6 + i * 0.85, 5.4, 0.7, p, fill=C_WHITE, title_size=13, border=C_LINE)

    rect(slide, 6.8, 1.35, 6.0, 5.2, C_SOFT_T, radius=True)
    section_chip(slide, 7.05, 1.55, "外部记忆共享", C_EXTERNAL)
    add_text(
        slide.shapes.add_textbox(Inches(7.05), Inches(2.05), Inches(5.5), Inches(0.4)),
        "结构化共享给系统 / 小艺入口",
        size=15,
        bold=True,
        color=C_INK,
    )
    for i, p in enumerate(
        [
            "Intents Kit：shareIntent 动作 + 实体",
            "供小艺对话 / 搜索 / 建议智慧分发",
            "云侧事件意图 → 时效提醒与撤销",
            "有频次与体量限制，需按 Schema 接入",
        ]
    ):
        box(slide, 7.1, 2.6 + i * 0.85, 5.4, 0.7, p, fill=C_WHITE, title_size=13, border=C_LINE)

    # 5 Overview
    slide = blank(prs)
    title_bar(slide, "总览架构：三类记忆与两条流向", "内部闭环保体验；外部共享保触达；小艺记忆是系统侧资产", 5)

    rect(slide, 0.5, 1.3, 4.0, 5.3, C_SOFT, radius=True)
    label(slide, 0.6, 1.4, 3.8, 0.3, "应用域", 13, True, C_PRIMARY)
    box(slide, 0.75, 1.85, 3.5, 0.85, "业务 UI / UIAbility", "产生行为与实体", fill=C_WHITE, title_size=13)
    box(
        slide,
        0.75,
        2.9,
        3.5,
        1.1,
        "内部 Memory Center",
        "Working / Episodic\nSemantic / Procedural",
        fill=C_SOFT_B,
        title_color=C_INTERNAL,
        title_size=13,
    )
    box(
        slide,
        0.75,
        4.2,
        3.5,
        0.85,
        "Intent Sharer",
        "动作完成时 / 将来时",
        fill=C_SOFT_T,
        title_color=C_EXTERNAL,
        title_size=13,
    )
    box(slide, 0.75, 5.25, 3.5, 0.95, "InsightIntentExecutor", "响应系统意图调用", fill=C_WHITE, title_size=13)

    rect(slide, 5.0, 1.3, 4.0, 5.3, C_SOFT_T, radius=True)
    label(slide, 5.1, 1.4, 3.8, 0.3, "系统智慧层", 13, True, C_EXTERNAL)
    box(slide, 5.25, 1.85, 3.5, 0.95, "Intents Kit / 意图框架", "索引 · 规律学习 · 分发", fill=C_WHITE, title_size=13)
    box(slide, 5.25, 3.05, 3.5, 0.85, "小艺对话 / 搜索 / 建议", "系统入口", fill=C_WHITE, title_size=13)
    box(
        slide,
        5.25,
        4.15,
        3.5,
        1.0,
        "小艺记忆 / 帮记",
        "系统侧长期记忆资产",
        fill=C_SOFT_M,
        title_color=C_MEMORY,
        title_size=13,
    )
    box(slide, 5.25, 5.4, 3.5, 0.85, "记忆授权应用", "系统 → 应用（反向）", fill=C_WHITE, title_size=13)

    rect(slide, 9.5, 1.3, 3.4, 5.3, C_SOFT_M, radius=True)
    label(slide, 9.6, 1.4, 3.2, 0.3, "智能体域", 13, True, C_MEMORY)
    box(slide, 9.7, 1.85, 3.0, 0.95, "小艺开放平台", "Agent 编排发布", fill=C_WHITE, title_size=13)
    box(
        slide,
        9.7,
        3.05,
        3.0,
        1.2,
        "长期记忆节点",
        "对话抽取 · Prompt 召回\nWorkflow 调用",
        fill=C_WHITE,
        title_size=13,
    )
    box(slide, 9.7, 4.5, 3.0, 0.95, "Skill / 意图工具化", "对接业务能力", fill=C_WHITE, title_size=13)
    box(
        slide,
        9.7,
        5.7,
        3.0,
        0.55,
        "属内部记忆范畴",
        fill=RGBColor(0xF8, 0xF0, 0xE8),
        title_size=12,
        title_color=C_MEMORY,
    )

    arrow_right(slide, 4.55, 4.5, 0.4, 0.22, C_EXTERNAL)
    label(slide, 4.35, 4.75, 0.8, 0.35, "外部共享", 9, False, C_EXTERNAL)
    arrow_right(slide, 9.05, 3.35, 0.4, 0.22, C_MEMORY)
    label(slide, 8.85, 3.6, 0.9, 0.35, "Agent记忆", 9, False, C_MEMORY)

    # 6 Scheme A
    slide = blank(prs)
    title_bar(slide, "方案 A｜外部记忆共享：Intents Kit", "把「可被系统理解的行为与实体」共享出去，换取小艺入口触达", 6)
    section_chip(slide, 0.6, 1.2, "外部共享", C_EXTERNAL)
    nodes = [
        (0.5, "应用业务事件", "播放/下单/出行…"),
        (3.0, "结构化 Intent", "Action + Entity"),
        (5.5, "shareIntent", "端侧 / 云侧"),
        (8.0, "意图框架", "索引 · 学习"),
        (10.5, "小艺入口", "对话/搜索/建议"),
    ]
    for i, (x, t, s) in enumerate(nodes):
        box(
            slide,
            x,
            1.75,
            2.2,
            1.15,
            t,
            s,
            fill=C_SOFT_T if i >= 3 else C_WHITE,
            title_size=14,
            border=C_PRIMARY if i == 2 else C_LINE,
        )
        if i < 4:
            arrow_right(slide, x + 2.25, 2.2, 0.35, 0.2, C_PRIMARY)

    rect(slide, 0.5, 3.3, 6.1, 3.25, C_SOFT, radius=True)
    label(slide, 0.7, 3.45, 5.5, 0.3, "共享侧要点", 14, True, C_PRIMARY)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(5.5), Inches(2.4))
    lines = [
        "完成时：已发生行为 → 搜索 + 建议",
        "将来时：预测意图 → 偏搜索",
        "辅助实体：位置等增强场景推荐",
        "限制：约 20 次/天 · 单次 ≤50KB",
        "建议积攒后批量共享，控制频次",
    ]
    tf = add_text(tb, lines[0], size=13, color=C_INK)
    for d in lines[1:]:
        add_para(tf, d, size=13, color=C_INK, space_before=8)

    rect(slide, 6.85, 3.3, 5.95, 3.25, C_SOFT_B, radius=True)
    label(slide, 7.05, 3.45, 5.5, 0.3, "调用侧要点", 14, True, C_INTERNAL)
    tb = slide.shapes.add_textbox(Inches(7.15), Inches(3.9), Inches(5.4), Inches(2.4))
    lines = [
        "insight_intent.json 注册意图",
        "实现 InsightIntentExecutor",
        "小艺触发 → 拉起 Ability / 落地页",
        "可绑定 UIAbility / Extension / 卡片",
        "垂域 Schema + 平台验收（常见）",
    ]
    tf = add_text(tb, lines[0], size=13, color=C_INK)
    for d in lines[1:]:
        add_para(tf, d, size=13, color=C_INK, space_before=8)

    # 7 Scheme A detail
    slide = blank(prs)
    title_bar(slide, "方案 A｜详细数据流（共享 ↔ 调用）", "外部共享是「推数据给系统」；意图调用是「系统拉起应用」", 7)
    rect(slide, 0.5, 1.25, 6.0, 5.35, C_SOFT_T, radius=True)
    label(slide, 0.7, 1.4, 5.5, 0.3, "路径 ① 意图共享（App → System）", 14, True, C_EXTERNAL)
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.3))
    flow1 = [
        "1. 业务层产生记忆片段（观看记录、航班、待办）",
        "2. 映射为 InsightIntent（name/version/id）",
        "3. 填充 IntentActionInfo + IntentEntityInfo",
        "4. insightIntent.shareIntent(context, intents)",
        "5. 系统构建本地索引 / 学习行为规律",
        "6. 合适时机在小艺建议/搜索露出",
    ]
    tf = add_text(tb, flow1[0], size=13, color=C_INK)
    for d in flow1[1:]:
        add_para(tf, d, size=13, color=C_INK, space_before=12)

    rect(slide, 6.85, 1.25, 5.95, 5.35, C_SOFT_B, radius=True)
    label(slide, 7.05, 1.4, 5.5, 0.3, "路径 ② 意图调用（System → App）", 14, True, C_INTERNAL)
    tb = slide.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.5), Inches(4.3))
    flow2 = [
        "1. 用户语音/点击建议卡片触发意图",
        "2. 小艺 NLU / 意图框架解析参数",
        "3. 匹配应用声明的 insightIntent API",
        "4. 拉起绑定组件（前台/后台/Extension）",
        "5. Executor 回调执行业务落地",
        "6. 返回 IntentResult，闭环体验",
    ]
    tf = add_text(tb, flow2[0], size=13, color=C_INK)
    for d in flow2[1:]:
        add_para(tf, d, size=13, color=C_INK, space_before=12)

    # 8 Scheme B
    slide = blank(prs)
    title_bar(slide, "方案 B｜内部记忆管理：应用 / Agent 长期记忆", "记忆服务本应用智能体验；不默认进入系统小艺记忆库", 8)
    section_chip(slide, 0.6, 1.2, "内部管理", C_INTERNAL)
    mems = [
        ("Working Memory", "当前会话上下文", "多轮对话窗口、临时状态", C_SOFT_B),
        ("Episodic / Fact", "情景与事实记忆", "历史事件、用户偏好事实", C_SOFT),
        ("Procedural", "程序/技能记忆", "常用操作路径、工作流", C_SOFT_T),
    ]
    for i, (t, s, d, f) in enumerate(mems):
        x = 0.55 + i * 4.2
        rect(slide, x, 1.7, 4.0, 2.0, f, radius=True, line=C_INTERNAL if i == 0 else C_LINE)
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.85), Inches(3.6), Inches(0.35)),
            t,
            15,
            True,
            C_INTERNAL,
        )
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.3), Inches(3.6), Inches(0.3)),
            s,
            13,
            True,
            C_INK,
        )
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.75), Inches(3.6), Inches(0.7)),
            d,
            12,
            False,
            C_MUTED,
        )

    rect(slide, 0.55, 4.05, 6.0, 2.5, C_WHITE, line=C_INTERNAL, radius=True)
    add_text(
        slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(5.5), Inches(0.35)),
        "B1 应用自建 Memory Center",
        15,
        True,
        C_INTERNAL,
    )
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(4.65), Inches(5.5), Inches(1.7))
    tf = add_text(tb, "· 本地存储 + 向量检索 / 结构化查询", 13)
    add_para(tf, "· 统一 save / load / query 接口", 13, space_before=6)
    add_para(tf, "· 可选跨端同步（分布式/软总线）", 13, space_before=6)
    add_para(tf, "· 输出：可再映射为 Intent 做外部共享", 13, space_before=6)

    rect(slide, 6.8, 4.05, 5.95, 2.5, C_WHITE, line=C_MEMORY, radius=True)
    add_text(
        slide.shapes.add_textbox(Inches(7.0), Inches(4.2), Inches(5.5), Inches(0.35)),
        "B2 小艺开放平台长期记忆",
        15,
        True,
        C_MEMORY,
    )
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(4.65), Inches(5.5), Inches(1.7))
    tf = add_text(tb, "· 编排开启长期记忆，对话自动抽取", 13)
    add_para(tf, "· Prompt 自动拼接 / Workflow 召回", 13, space_before=6)
    add_para(tf, "· 适合 AI 助理、陪伴类智能体", 13, space_before=6)
    add_para(tf, "· 仍是智能体域记忆，≠ 系统小艺记忆库", 13, color=C_WARN, space_before=6)

    # 9 Scheme C
    slide = blank(prs)
    title_bar(slide, "方案 C｜小艺记忆 / 帮记（系统侧记忆）", "系统级记忆资产；三方应用通常「被授权读取」而非「直接写入」", 9)
    section_chip(slide, 0.6, 1.2, "系统记忆", C_MEMORY)
    box(
        slide,
        0.5,
        1.7,
        3.5,
        1.5,
        "写入来源（系统侧）",
        "用户口述「帮我记住…」\n拖拽文本/图片/文档\n系统场景沉淀",
        fill=C_WHITE,
        title_size=13,
        border=C_LINE,
    )
    arrow_right(slide, 4.05, 2.25, 0.4, 0.2, C_MEMORY)
    box(
        slide,
        4.5,
        1.7,
        4.3,
        1.5,
        "小艺记忆 / 小艺帮记",
        "卡证 · 地址 · 重要日 · 订单\n待办 · 合集 · 知识收藏",
        fill=C_SOFT_M,
        title_color=C_MEMORY,
        title_size=16,
        border=C_MEMORY,
    )
    arrow_right(slide, 8.9, 2.25, 0.4, 0.2, C_MEMORY)
    box(
        slide,
        9.4,
        1.7,
        3.4,
        1.5,
        "记忆授权应用",
        "用户授权后，应用可使用\n记忆增强体验（如输入法）",
        fill=C_WHITE,
        title_size=13,
        border=C_MEMORY,
    )

    cards = [
        ("能做什么", "系统理解用户长期偏好；跨场景召回；增强小艺回答与建议"),
        (
            "对三方的现实约束",
            "无公开「任意写入小艺记忆库」的通用开发接口；写入以系统/用户操作为主",
        ),
        (
            "应用可配合的方式",
            "① 意图共享让小艺更懂业务\n② 申请成为记忆授权应用（反向读）\n③ 自建内部记忆 + Agent 长期记忆",
        ),
    ]
    for i, (t, d) in enumerate(cards):
        x = 0.5 + i * 4.2
        rect(slide, x, 3.6, 4.0, 2.9, C_SOFT_M if i == 0 else C_SOFT, radius=True)
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(3.6), Inches(0.4)),
            t,
            15,
            True,
            C_MEMORY if i == 0 else C_INK,
        )
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.35), Inches(3.6), Inches(1.9)),
            d,
            13,
            False,
            C_INK,
        )

    # 10 Scheme D
    slide = blank(prs)
    title_bar(slide, "方案 D｜双轨组合架构（推荐）", "内部记忆保「越用越懂」；外部共享保「小艺能找到你」；系统记忆作增强可选", 10)
    lanes = [
        (1.25, "采集层", "UI 行为 / 业务事件 / 对话内容", C_SOFT),
        (2.35, "内部记忆层", "Memory Center +（可选）Agent 长期记忆", C_SOFT_B),
        (3.45, "对外共享层", "Intent 映射 · shareIntent · 云侧事件", C_SOFT_T),
        (4.55, "系统分发层", "Intents Kit · 小艺对话/搜索/建议", C_SOFT),
        (5.65, "系统记忆层（可选）", "小艺记忆/帮记 · 记忆授权回灌应用", C_SOFT_M),
    ]
    for y, t, d, f in lanes:
        rect(slide, 0.55, y, 2.6, 0.95, f, radius=True)
        add_text(
            slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(2.3), Inches(0.3)),
            t,
            13,
            True,
            C_INK,
        )
        add_text(
            slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.5), Inches(2.3), Inches(0.35)),
            d,
            10,
            False,
            C_MUTED,
        )
        if y < 5.65:
            arrow_down(slide, 1.7, y + 0.95, 0.2, 0.14, C_PRIMARY)

    rect(slide, 3.5, 1.25, 9.2, 5.35, C_WHITE, line=C_LINE, radius=True)
    add_text(
        slide.shapes.add_textbox(Inches(3.75), Inches(1.45), Inches(8.7), Inches(0.4)),
        "组合原则",
        16,
        True,
        C_PRIMARY,
    )
    principles = [
        ("1. 先内后外", "先把应用内记忆治理清楚（结构、时效、删除、授权），再挑选「值得共享」的子集映射为 Intent。"),
        ("2. 共享最小化", "只共享系统入口真正需要的动作/实体，避免把完整画像无差别外抛。"),
        ("3. 双向闭环", "共享解决触达；Executor 解决落地；内部记忆解决个性化连续性。"),
        ("4. 系统记忆定位", "小艺记忆/帮记作为系统资产与用户授权增强，不作为三方主写入通道。"),
        ("5. 合规与频控", "遵守 shareIntent 频次/体量限制；云侧事件要有时效与撤销。"),
    ]
    tb = slide.shapes.add_textbox(Inches(3.8), Inches(2.0), Inches(8.6), Inches(4.3))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (a, b) in enumerate(principles):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = a + "  "
        set_run(r1, 13, True, C_PRIMARY)
        r2 = p.add_run()
        r2.text = b
        set_run(r2, 13, False, C_INK)

    # 11 Compare
    slide = blank(prs)
    title_bar(slide, "方案对比与选型建议", "按目标选主路径，而不是一套方案打天下", 11)
    headers = ["维度", "A 外部意图共享", "B 内部长期记忆", "C 小艺记忆/帮记"]
    widths = [2.0, 3.4, 3.4, 3.4]
    x0 = 0.55
    rect(slide, x0, 1.3, 12.2, 0.5, C_DARK)
    xx = x0
    for i, h in enumerate(headers):
        add_text(
            slide.shapes.add_textbox(Inches(xx), Inches(1.38), Inches(widths[i]), Inches(0.35)),
            h,
            12,
            True,
            C_WHITE,
            PP_ALIGN.CENTER,
        )
        xx += widths[i]

    rows = [
        ["记忆权属", "系统索引/分发侧", "应用或 Agent 域", "系统侧用户记忆"],
        ["典型能力", "对话/搜索/建议触达", "个性化连续对话", "跨场景个人记忆"],
        ["三方可写性", "可（受限共享）", "完全可写", "基本不可直接写"],
        ["与小艺关系", "强（主对接通道）", "中（智能体内）", "强但偏系统能力"],
        ["推荐场景", "要被小艺叫到/推荐", "应用内 AI 助理", "用户级记忆增强"],
    ]
    for ri, row in enumerate(rows):
        y = 1.85 + ri * 0.7
        bg = C_SOFT if ri % 2 == 0 else C_WHITE
        rect(slide, x0, y, 12.2, 0.65, bg)
        xx = x0
        for ci, cell in enumerate(row):
            add_text(
                slide.shapes.add_textbox(
                    Inches(xx + 0.08), Inches(y + 0.15), Inches(widths[ci] - 0.1), Inches(0.4)
                ),
                cell,
                12,
                ci == 0,
                C_INK if ci == 0 else C_MUTED,
                PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT,
            )
            xx += widths[ci]

    rect(slide, 0.55, 5.5, 12.2, 1.15, C_SOFT_T, radius=True)
    add_text(
        slide.shapes.add_textbox(Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.35)),
        "落地建议：默认走方案 D（B 内部治理 + A 外部共享）；C 作为系统增强与授权回灌，不作为主写入通道。",
        14,
        True,
        C_PRIMARY,
    )
    add_text(
        slide.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.35)),
        "若只有一个目标——让小艺能调用应用能力：优先 A；若只有个性化对话：优先 B；若讨论系统记忆产品能力：看 C。",
        12,
        False,
        C_INK,
    )

    # 12 Roadmap
    slide = blank(prs)
    title_bar(slide, "落地步骤（建议节奏）", "从内部治理到外部共享，再到可选的系统记忆增强", 12)
    steps = [
        ("Step 1", "盘点记忆资产", "哪些是内部个性化，哪些适合对外共享", C_INTERNAL),
        ("Step 2", "建设内部 Memory", "分层存储、检索、删除与授权", C_INTERNAL),
        ("Step 3", "注册并映射 Intent", "insight_intent.json + Schema", C_EXTERNAL),
        ("Step 4", "实现共享与调用", "shareIntent + Executor 闭环", C_EXTERNAL),
        ("Step 5", "评估系统记忆增强", "帮记协同 / 记忆授权（可选）", C_MEMORY),
    ]
    for i, (n, t, d, c) in enumerate(steps):
        x = 0.45 + i * 2.55
        rect(slide, x, 1.5, 2.4, 3.2, C_WHITE, line=c, radius=True)
        chip = rect(slide, x + 0.2, 1.7, 1.3, 0.32, c, radius=True)
        add_text(chip, n, 12, True, C_WHITE, PP_ALIGN.CENTER)
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.3), Inches(2.1), Inches(0.8)),
            t,
            15,
            True,
            C_INK,
            PP_ALIGN.CENTER,
        )
        add_text(
            slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.2), Inches(2.1), Inches(1.2)),
            d,
            12,
            False,
            C_MUTED,
            PP_ALIGN.CENTER,
        )
        if i < 4:
            arrow_right(slide, x + 2.4, 2.9, 0.18, 0.18, C_LINE)

    rect(slide, 0.55, 5.1, 12.2, 1.5, C_DARK, radius=True)
    add_text(
        slide.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.4)),
        "一句话结论",
        14,
        True,
        C_PRIMARY,
    )
    add_text(
        slide.shapes.add_textbox(Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.5)),
        "应用要「给小艺用记忆」：内部可完整管理；对外走意图共享；小艺记忆/帮记是系统资产，宜协同不宜硬写。",
        16,
        True,
        C_WHITE,
    )

    out_paths = [
        "/workspace/artifacts/应用记忆对接小艺_方案架构.pptx",
        "/opt/cursor/artifacts/应用记忆对接小艺_方案架构.pptx",
    ]
    for path in out_paths:
        prs.save(path)
        print("saved", path)
    print("slides", len(prs.slides))


if __name__ == "__main__":
    build()
