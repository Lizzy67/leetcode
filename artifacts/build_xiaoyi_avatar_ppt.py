#!/usr/bin/env python3
"""精简版：小艺应用分身平台架构 PPT（约 7 页）。"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree

C_DARK = RGBColor(0x12, 0x1A, 0x22)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_INK = RGBColor(0x1A, 0x2B, 0x3C)
C_MUTED = RGBColor(0x5A, 0x6B, 0x7C)
C_TEAL = RGBColor(0x0E, 0x7C, 0x86)
C_TEAL_D = RGBColor(0x0A, 0x5C, 0x64)
C_BLUE = RGBColor(0x2F, 0x6F, 0xE5)
C_ORANGE = RGBColor(0xC4, 0x5C, 0x26)
C_LINE = RGBColor(0xD0, 0xD8, 0xE0)
C_SOFT = RGBColor(0xE8, 0xF0, 0xF3)
C_SOFT_B = RGBColor(0xE3, 0xEC, 0xFB)
C_SOFT_T = RGBColor(0xD9, 0xEF, 0xF1)
C_SOFT_O = RGBColor(0xFA, 0xEB, 0xE0)
C_WARN = RGBColor(0xB5, 0x6E, 0x1A)
C_OK = RGBColor(0x1F, 0x8A, 0x5B)

TOTAL = 7
FONT = "Microsoft YaHei"


def set_run(run, size=18, bold=False, color=C_INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = etree.SubElement(
            r_pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea"
        )
    ea.set("typeface", FONT)


def add_text(shape, text, size=18, bold=False, color=C_INK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return tf


def add_para(tf, text, size=13, bold=False, color=C_INK, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)


def rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.15)
    if radius:
        s.adjustments[0] = 0.08
    return s


def box(slide, x, y, w, h, title, sub=None, fill=C_SOFT, tc=C_INK, sc=C_MUTED, ts=13, ss=11, border=None):
    s = rect(slide, x, y, w, h, fill, line=border)
    tf = s.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(6)
    r = tf.paragraphs[0].add_run()
    r.text = title
    set_run(r, ts, True, tc)
    if sub:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(3)
        rr = p.add_run()
        rr.text = sub
        set_run(rr, ss, False, sc)
    return s


def arrow_r(slide, x, y, w=0.32, h=0.18, color=C_TEAL):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def arrow_d(slide, x, y, w=0.18, h=0.28, color=C_TEAL):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def chip(slide, x, y, text, fill, w=1.35):
    s = rect(slide, x, y, w, 0.28, fill)
    add_text(s, text, 11, True, C_WHITE, PP_ALIGN.CENTER)


def title_bar(slide, title, sub, page):
    rect(slide, 0, 0, 13.333, 0.07, C_TEAL, radius=False)
    add_text(
        slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(11), Inches(0.42)),
        title, 26, True, C_INK,
    )
    add_text(
        slide.shapes.add_textbox(Inches(0.55), Inches(0.68), Inches(11.5), Inches(0.3)),
        sub, 13, False, C_MUTED,
    )
    add_text(
        slide.shapes.add_textbox(Inches(12.1), Inches(7.1), Inches(0.9), Inches(0.28)),
        f"{page}/{TOTAL}", 11, False, C_MUTED, PP_ALIGN.RIGHT,
    )
    rect(slide, 0.55, 7.0, 12.2, 0.01, C_LINE, radius=False)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.333, 7.5, C_WHITE, radius=False)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== 1 Cover =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.333, 7.5, C_DARK, radius=False)
    rect(s, 0, 0, 0.16, 7.5, C_TEAL, radius=False)
    rect(s, 0, 6.95, 13.333, 0.55, C_TEAL_D, radius=False)
    add_text(s.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(11), Inches(0.7)),
             "小艺应用分身平台", 40, True, C_WHITE)
    add_text(s.shapes.add_textbox(Inches(0.85), Inches(2.85), Inches(11), Inches(0.45)),
             "App Avatar · Skill Tag Channel · 隔离接入 · 记忆与习惯数据", 18, False, RGBColor(0xB0, 0xC4, 0xCC))
    add_text(s.shapes.add_textbox(Inches(0.85), Inches(3.6), Inches(11), Inches(0.35)),
             "图库分身 · 浏览器分身 · 备忘录分身 · 地图分身 …", 14, False, C_TEAL)
    add_text(s.shapes.add_textbox(Inches(0.85), Inches(7.1), Inches(10), Inches(0.28)),
             "配置一体化 · 能力平台化 · 服务/设备隔离", 12, False, RGBColor(0xD0, 0xE8, 0xEA))

    # ===== 2 分身模型 =====
    s = blank(prs)
    title_bar(s, "01  分身是什么", "每个应用分身 = 独立会话运行时 + 一套可配置资产包", 2)

    # left: identity
    rect(s, 0.5, 1.2, 4.0, 5.4, C_SOFT_T)
    chip(s, 0.7, 1.4, "分身身份", C_TEAL, 1.4)
    assets = [
        ("system prompt", "人设与边界"),
        ("soul.md", "价值观 / 语气"),
        ("agent.md", "目标与策略"),
        ("userProfile.md", "用户画像视图"),
        ("memory.md", "长期记忆"),
        ("Skills / Tools.md", "可调用能力"),
        ("workspace", "会话工作区"),
    ]
    for i, (a, b) in enumerate(assets):
        y = 1.9 + i * 0.58
        box(s, 0.7, y, 3.6, 0.5, f"{a}  ·  {b}", fill=C_WHITE, ts=12, border=C_LINE)

    # center arrow concept
    arrow_r(s, 4.65, 3.6, 0.4, 0.22)

    # right: runtime properties
    rect(s, 5.2, 1.2, 7.6, 5.4, C_SOFT)
    chip(s, 5.45, 1.4, "运行时约束", C_BLUE, 1.5)
    props = [
        ("独立会话", "分身间上下文隔离，不串话、不串记忆默认视图"),
        ("按应用隔离", "图库 / 浏览器 / 备忘录 / 地图 … 各自 namespace"),
        ("按设备隔离", "Phone / PC / 车机 … profile 与能力集可不同"),
        ("配置驱动", "一键接入：声明分身 ID + 资产包 + 接口 + Skill 绑定"),
        ("启动可观测", "习惯数据拉取有超时预算，失败降级不阻塞对话"),
    ]
    for i, (t, d) in enumerate(props):
        y = 1.95 + i * 0.85
        rect(s, 5.45, y, 7.1, 0.72, C_WHITE, line=C_LINE)
        add_text(s.shapes.add_textbox(Inches(5.65), Inches(y + 0.08), Inches(6.7), Inches(0.28)),
                 t, 14, True, C_INK)
        add_text(s.shapes.add_textbox(Inches(5.65), Inches(y + 0.36), Inches(6.7), Inches(0.28)),
                 d, 12, False, C_MUTED)

    # ===== 3 Skill Tag Channel =====
    s = blank(prs)
    title_bar(s, "02  Skill Tag Channel", "上架打 Tag 划领域；运行时按 Tag 装配 Skill / Memory / Soul，再开对话", 3)

    # pipeline
    steps = [
        (0.5, "Skill 上架", "domain tags\ncapability tags"),
        (3.0, "Channel 绑定", "avatar ↔ tag set\n+ 私有 skill"),
        (5.5, "SkillSearch", "按 tag 召回\n排序 / 裁剪"),
        (8.0, "运行时装配", "Soul · Memory\nTools · Prompt"),
        (10.5, "独立会话", "分身对话\n可回流记忆"),
    ]
    for i, (x, t, sub) in enumerate(steps):
        box(s, x, 1.25, 2.25, 1.35, t, sub, fill=C_SOFT_T if i in (1, 3) else C_WHITE,
            tc=C_TEAL, ts=14, ss=11, border=C_TEAL if i == 2 else C_LINE)
        if i < 4:
            arrow_r(s, x + 2.28, 1.8, 0.3, 0.18)

    # two columns detail
    rect(s, 0.5, 2.95, 6.0, 3.65, C_SOFT_B)
    chip(s, 0.7, 3.15, "Tag 体系（建议）", C_BLUE, 1.9)
    tb = s.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(5.5), Inches(2.8))
    tf = add_text(tb, "domain.*     gallery / browser / notes / map / …", 13, True, C_INK)
    add_para(tf, "device.*      phone / pc / car / pad", 13, True, C_INK, 10)
    add_para(tf, "cap.*         search / edit / summarize / navigate", 13, True, C_INK, 10)
    add_para(tf, "risk.*        local-only / needs-network / PII", 13, True, C_INK, 10)
    add_para(tf, " ", 8, False, C_MUTED, 4)
    add_para(tf, "示例：小艺图库分身 Channel =", 12, False, C_MUTED, 4)
    add_para(tf, "{domain.gallery, device.phone} ∪ 私有 Skills", 13, True, C_TEAL, 4)

    rect(s, 6.8, 2.95, 6.0, 3.65, C_SOFT_T)
    chip(s, 7.0, 3.15, "运行时加载规则", C_TEAL, 1.9)
    tb = s.shapes.add_textbox(Inches(7.05), Inches(3.6), Inches(5.5), Inches(2.8))
    tf = add_text(tb, "1. 解析 AvatarConfig → required_tags / deny_tags", 13)
    add_para(tf, "2. SkillSearch(tags) → TopK（受配额限制）", 13, space_before=8)
    add_para(tf, "3. 合并：私有 Skill 优先于公共 Skill", 13, space_before=8)
    add_para(tf, "4. 同步装载 soul / memory / tools 视图", 13, space_before=8)
    add_para(tf, "5. 开独立 Session（session_id 含 avatar+device）", 13, space_before=8)

    # ===== 4 隔离 + 一键接入 =====
    s = blank(prs)
    title_bar(s, "03  隔离模型 × 一键接入", "隔离保证安全边界；配置平台保证新应用接入简单", 4)

    # isolation matrix
    rect(s, 0.5, 1.2, 6.3, 5.4, C_SOFT)
    chip(s, 0.7, 1.4, "二维隔离", C_ORANGE, 1.4)
    # matrix header
    headers = ["", "Phone", "PC", "车机"]
    apps = ["图库", "浏览器", "备忘录", "地图"]
    for i, h in enumerate(headers):
        x = 1.0 + i * 1.35
        box(s, x, 2.0, 1.25, 0.45, h or "应用\\设备", fill=C_DARK, tc=C_WHITE, ts=11)
    for r, app in enumerate(apps):
        box(s, 1.0, 2.55 + r * 0.7, 1.25, 0.55, app, fill=C_SOFT_O, tc=C_ORANGE, ts=12)
        for c in range(3):
            box(s, 2.35 + c * 1.35, 2.55 + r * 0.7, 1.25, 0.55, "独立\n分身实例",
                fill=C_WHITE, ts=10, ss=9, border=C_LINE)
    add_text(s.shapes.add_textbox(Inches(0.75), Inches(5.55), Inches(5.8), Inches(0.8)),
             "隔离键建议：tenant × app_id × device_class × avatar_id\n会话 / 记忆 / Skill 可见性均按此键切分", 12, False, C_MUTED)

    # onboarding
    rect(s, 7.05, 1.2, 5.75, 5.4, C_SOFT_T)
    chip(s, 7.25, 1.4, "一键接入（配置一体化）", C_TEAL, 2.6)
    steps2 = [
        "① 注册 App / 分身（ID、名称、设备范围）",
        "② 上传或引用资产包（prompt/soul/agent/…）",
        "③ 绑定 Channel Tags + 私有/公共 Skills",
        "④ 配置 Habit API（受统一 Schema 约束）",
        "⑤ 发布：灰度 → 全量；可一键回滚",
    ]
    tb = s.shapes.add_textbox(Inches(7.25), Inches(2.0), Inches(5.3), Inches(3.2))
    tf = add_text(tb, steps2[0], 14, True, C_INK)
    for line in steps2[1:]:
        add_para(tf, line, 14, True, C_INK, 14)
    add_text(s.shapes.add_textbox(Inches(7.25), Inches(5.5), Inches(5.3), Inches(0.8)),
             "平台化目标：接入只改配置，不改内核；\n能力复用靠 Tag/公共 Skill，差异靠私有包。", 12, False, C_TEAL)

    # ===== 5 Memory + Habit API =====
    s = blank(prs)
    title_bar(s, "04  记忆回流 / 在线提取 / 习惯数据", "分身记忆可写回；启动时按 Schema 拉取习惯，且必须有时延预算", 5)

    # three memory modes
    modes = [
        ("数据回流", "user.md / memory.md\n会话结束后结构化写回", C_SOFT_B, C_BLUE),
        ("在线提取", "对话摘要 / 偏好增量\n低延迟增量合并", C_SOFT_T, C_TEAL),
        ("习惯查询", "启动时调 App Habit API\n只取启动必要上下文", C_SOFT_O, C_ORANGE),
    ]
    for i, (t, d, f, c) in enumerate(modes):
        x = 0.5 + i * 4.2
        rect(s, x, 1.2, 4.0, 1.55, f, line=c)
        chip(s, x + 0.2, 1.35, t, c, 1.35)
        add_text(s.shapes.add_textbox(Inches(x + 0.2), Inches(1.8), Inches(3.6), Inches(0.8)),
                 d, 13, False, C_INK)

    # schema + perf
    rect(s, 0.5, 3.0, 7.6, 3.6, C_SOFT)
    chip(s, 0.7, 3.2, "Habit API Schema（参考）", C_ORANGE, 2.5)
    schema = (
        "GET /v1/avatars/{avatar_id}/habits?device={device}&limit=N\n"
        "Response {\n"
        "  schema_version: \"1.0\",\n"
        "  app_id, user_id_hash, device_class,\n"
        "  fetched_at, ttl_sec,\n"
        "  habits: [{ key, value, score?, updated_at, source }],\n"
        "  recent_actions: [{ action, entity_type?, entity_id?, ts }],\n"
        "  hints: [{ type, text, priority }]   // 可选，≤3\n"
        "}"
    )
    add_text(s.shapes.add_textbox(Inches(0.75), Inches(3.65), Inches(7.1), Inches(2.7)),
             schema, 11, False, C_INK)

    rect(s, 8.35, 3.0, 4.45, 3.6, C_SOFT_O, line=C_ORANGE)
    chip(s, 8.55, 3.2, "启动性能约束", C_ORANGE, 1.8)
    tb = s.shapes.add_textbox(Inches(8.55), Inches(3.7), Inches(4.0), Inches(2.7))
    tf = add_text(tb, "超时预算：≤ 150–300ms", 13, True, C_INK)
    add_para(tf, "超时/失败：降级空习惯，不阻塞开聊", 12, False, C_MUTED, 8)
    add_para(tf, "payload ≤ 8–16KB；habits ≤ 20 条", 12, False, C_MUTED, 8)
    add_para(tf, "recent_actions ≤ 10；禁止大段原文", 12, False, C_MUTED, 8)
    add_para(tf, "可缓存：ETag / ttl_sec（建议 ≥60s）", 12, False, C_MUTED, 8)
    add_para(tf, "仅启动必要字段；详情懒加载", 12, False, C_MUTED, 8)
    add_para(tf, "P99 纳入分身启动 SLO 看板", 12, True, C_WARN, 10)

    # ===== 6 Public Skill boundaries =====
    s = blank(prs)
    title_bar(s, "05  公共 Skill 边界", "公共可复用，但不能越权；私有承载业务差异与数据敏感操作", 6)

    # principles
    principles = [
        ("最小权限", "公共 Skill 默认无 App 私有数据访问；需显式授权 scope"),
        ("领域中立", "不绑定单一应用业务语义；用 Tag 表达适用域，而非硬编码 app_id"),
        ("副作用可控", "写操作（删除/支付/发帖）不得进公共；或强制二次确认 + 私有封装"),
        ("可组合", "输入输出 Schema 稳定；可被多 Channel 引用，版本可灰度"),
        ("可撤销", "公共 Skill 下架/降级不影响私有核心链路；Channel 有 deny_tags"),
        ("设备感知", "声明 device 兼容性；车机等受限环境自动过滤高风险能力"),
    ]
    for i, (t, d) in enumerate(principles):
        col, row = i % 2, i // 2
        x = 0.5 + col * 6.4
        y = 1.2 + row * 1.35
        rect(s, x, y, 6.15, 1.2, C_SOFT if col == 0 else C_SOFT_T, line=C_LINE)
        add_text(s.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(5.7), Inches(0.3)),
                 t, 15, True, C_TEAL)
        add_text(s.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.55), Inches(5.7), Inches(0.5)),
                 d, 12, False, C_INK)

    # ===== 7 Summary one-pager =====
    s = blank(prs)
    title_bar(s, "06  总览一页纸", "逻辑链：分身资产 → Tag Channel 装配 → 隔离运行 → 记忆/习惯闭环 → 公共边界", 7)

    # flow boxes compact
    nodes = [
        ("配置平台", "一键接入\n资产+Tag+API"),
        ("Avatar\nRuntime", "独立会话\n按隔离键"),
        ("SkillSearch", "Tag 召回\n私有优先"),
        ("Memory\nPlane", "回流/摘要\nHabit 拉取"),
        ("对话\n落地", "分身能力\n可观测"),
    ]
    for i, (t, d) in enumerate(nodes):
        x = 0.45 + i * 2.55
        box(s, x, 1.25, 2.35, 1.4, t, d, fill=C_SOFT_T if i % 2 == 0 else C_WHITE,
            tc=C_TEAL, ts=13, ss=11, border=C_TEAL)
        if i < 4:
            arrow_r(s, x + 2.38, 1.85, 0.2, 0.16)

    # decision bullets
    rect(s, 0.5, 3.0, 12.3, 3.55, C_DARK)
    add_text(s.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.35)),
             "拍板结论", 14, True, C_TEAL)
    bullets = [
        "1. 以「应用分身」为产品单元：独立会话 + 完整资产包（prompt/soul/agent/profile/memory/skills/tools/workspace）。",
        "2. 以 Skill Tag Channel 做领域装配：上架打标签，运行时 SkillSearch；私有 Skill 优先，公共 Skill 可选。",
        "3. 隔离键 = 应用 × 设备（可扩展租户）；配置一体化支撑一键接入，能力平台化靠 Tag/公共池。",
        "4. 记忆三件套：回流（user/memory.md）+ 在线摘要 + 启动 Habit API（强 Schema + 超时降级）。",
        "5. 公共 Skill 边界：最小权限、领域中立、副作用可控、可组合、可撤销、设备感知。",
    ]
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.65), Inches(11.7), Inches(2.7))
    tf = add_text(tb, bullets[0], 13, False, C_WHITE)
    for b in bullets[1:]:
        add_para(tf, b, 13, False, C_WHITE, 8)

    paths = [
        "/workspace/artifacts/小艺应用分身_平台架构.pptx",
        "/opt/cursor/artifacts/小艺应用分身_平台架构.pptx",
    ]
    for p in paths:
        prs.save(p)
        print("saved", p)
    print("slides", len(prs.slides))


if __name__ == "__main__":
    build()
